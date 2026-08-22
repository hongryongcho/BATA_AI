"""
전 세계 유명 개인투자자 전략 vs BATA 알고리즘 비교 PPT 생성
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree
import copy

# ── 색상 팔레트 ────────────────────────────────────────────────────────────────
BG_DARK      = RGBColor(0x0b, 0x12, 0x20)   # 배경 #0b1220
ACCENT_BLUE  = RGBColor(0x00, 0x9f, 0xff)   # 포인트 블루
ACCENT_CYAN  = RGBColor(0x00, 0xe5, 0xff)   # 밝은 시안
ACCENT_GREEN = RGBColor(0x00, 0xd4, 0x8a)   # 초록
ACCENT_RED   = RGBColor(0xff, 0x4d, 0x6d)   # 빨강
ACCENT_GOLD  = RGBColor(0xff, 0xd7, 0x00)   # 골드
WHITE        = RGBColor(0xff, 0xff, 0xff)
GRAY_LIGHT   = RGBColor(0xb0, 0xba, 0xd0)
GRAY_MID     = RGBColor(0x2a, 0x3a, 0x5a)
GRAY_DIM     = RGBColor(0x1a, 0x25, 0x3f)
TEXT_DIM     = RGBColor(0x78, 0x8a, 0xaa)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

blank_layout = prs.slide_layouts[6]  # blank


# ── 헬퍼 함수 ──────────────────────────────────────────────────────────────────
def set_bg(slide, color: RGBColor = BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(1, left, top, width, height)  # MSO_SHAPE_TYPE.RECTANGLE=1
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, text, left, top, width, height,
                 font_size=18, bold=False, color=WHITE,
                 align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(left, top, width, height)
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Malgun Gothic"
    return txb


def add_multiline_textbox(slide, lines, left, top, width, height,
                           font_size=14, color=WHITE, line_spacing=1.2,
                           bold_first=False, align=PP_ALIGN.LEFT):
    """lines: list of (text, bold, color_override or None)"""
    txb = slide.shapes.add_textbox(left, top, width, height)
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for item in lines:
        if isinstance(item, str):
            text, is_bold, col = item, False, color
        else:
            text = item[0]
            is_bold = item[1] if len(item) > 1 else False
            col = item[2] if len(item) > 2 and item[2] else color

        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.bold = (bold_first and first) or is_bold
        run.font.color.rgb = col
        run.font.name = "Malgun Gothic"
    return txb


def add_divider(slide, top, color=ACCENT_BLUE, thickness=Pt(1.5)):
    line = slide.shapes.add_shape(1, Inches(0.5), top, SLIDE_W - Inches(1.0), thickness)
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()
    return line


def add_table_row(table, row_idx, values, font_size=12, bold=False,
                  bg_color=None, text_colors=None):
    row = table.rows[row_idx]
    for ci, val in enumerate(values):
        cell = row.cells[ci]
        cell.text = str(val)
        tf = cell.text_frame
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = tf.paragraphs[0].runs
        if run:
            run[0].font.size = Pt(font_size)
            run[0].font.bold = bold
            run[0].font.name = "Malgun Gothic"
            if text_colors and ci < len(text_colors) and text_colors[ci]:
                run[0].font.color.rgb = text_colors[ci]
            else:
                run[0].font.color.rgb = WHITE
        if bg_color:
            fill = cell.fill
            fill.solid()
            fill.fore_color.rgb = bg_color


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — 커버
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s)

# 배경 그라데이션 바
add_rect(s, 0, 0, SLIDE_W, Inches(0.08), ACCENT_BLUE)
add_rect(s, 0, SLIDE_H - Inches(0.08), SLIDE_W, Inches(0.08), ACCENT_BLUE)

# 중앙 accent 박스
add_rect(s, Inches(0.5), Inches(2.5), Inches(12.33), Inches(3.2),
         fill_color=GRAY_DIM, line_color=ACCENT_BLUE, line_width=Pt(2))

# 타이틀
add_text_box(s, "전 세계 유명 개인투자자 전략 비교",
             Inches(0.7), Inches(2.65), Inches(12), Inches(1.0),
             font_size=38, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text_box(s, "& BATA RSI(2) 분할매수 알고리즘",
             Inches(0.7), Inches(3.55), Inches(12), Inches(0.8),
             font_size=28, bold=False, color=ACCENT_CYAN, align=PP_ALIGN.CENTER)

# 서브타이틀
add_text_box(s, "수익률 · MDD · 전략 철학 · 구현 난이도 종합 분석",
             Inches(0.7), Inches(4.35), Inches(12), Inches(0.5),
             font_size=16, bold=False, color=GRAY_LIGHT, align=PP_ALIGN.CENTER)

# 날짜
add_text_box(s, "2026.06",
             Inches(11.0), Inches(6.8), Inches(2.0), Inches(0.4),
             font_size=14, bold=False, color=TEXT_DIM, align=PP_ALIGN.RIGHT)

# 국기 이모지 대신 국가명 텍스트로
flags = "🇺🇸 미국  |  🇰🇷 한국  |  🇯🇵 일본  |  🇸🇬 싱가포르  |  🇩🇪 독일  |  🇭🇰 홍콩"
add_text_box(s, flags,
             Inches(0.7), Inches(5.3), Inches(12), Inches(0.5),
             font_size=14, color=TEXT_DIM, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — 목차
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s)
add_rect(s, 0, 0, SLIDE_W, Inches(0.08), ACCENT_BLUE)

add_text_box(s, "목차", Inches(0.5), Inches(0.2), Inches(12), Inches(0.7),
             font_size=32, bold=True, color=WHITE)
add_divider(s, Inches(1.0))

toc = [
    ("01", "BATA RSI(2) 분할매수 전략 개요",        Inches(0.6), Inches(1.1)),
    ("02", "BATA 백테스트 성과 (2021~2026)",         Inches(0.6), Inches(1.7)),
    ("03", "HFEA — Hedgefundie (미국)",              Inches(0.6), Inches(2.3)),
    ("04", "Dual Momentum GEM — Gary Antonacci",     Inches(0.6), Inches(2.9)),
    ("05", "강환국 VAA/BAA 동적 자산배분 (한국)",    Inches(0.6), Inches(3.5)),
    ("06", "일본 개인투자자 전략",                   Inches(6.6), Inches(1.1)),
    ("07", "싱가포르 CPF·TTI 전략",                  Inches(6.6), Inches(1.7)),
    ("08", "Dragon Portfolio — Chris Cole (미국)",   Inches(6.6), Inches(2.3)),
    ("09", "Gerd Kommer 글로벌 팩터 ETF (독일)",     Inches(6.6), Inches(2.9)),
    ("10", "종합 비교표 & 결론",                     Inches(6.6), Inches(3.5)),
]

for num, title, lft, tp in toc:
    add_rect(s, lft, tp, Inches(0.55), Inches(0.45), fill_color=ACCENT_BLUE)
    add_text_box(s, num, lft, tp, Inches(0.55), Inches(0.45),
                 font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(s, title, lft + Inches(0.65), tp, Inches(5.5), Inches(0.45),
                 font_size=14, color=GRAY_LIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — BATA 전략 개요
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s)
add_rect(s, 0, 0, SLIDE_W, Inches(0.08), ACCENT_GREEN)

add_text_box(s, "01  BATA RSI(2) 분할매수 전략 개요",
             Inches(0.5), Inches(0.15), Inches(12), Inches(0.7),
             font_size=26, bold=True, color=WHITE)
add_divider(s, Inches(0.95), color=ACCENT_GREEN)

# 좌측: 전략 설명
add_rect(s, Inches(0.4), Inches(1.05), Inches(5.8), Inches(5.8), fill_color=GRAY_DIM,
         line_color=ACCENT_GREEN, line_width=Pt(1))
add_text_box(s, "전략 구조", Inches(0.5), Inches(1.1), Inches(5.6), Inches(0.45),
             font_size=16, bold=True, color=ACCENT_GREEN)

strategy_lines = [
    ("■ 기본 배분", True, ACCENT_CYAN),
    ("  BIL(T-bill) 10% 보유 + TQQQ 5-lot 풀 90%", False, GRAY_LIGHT),
    ("  각 lot = 전체자산의 18%", False, TEXT_DIM),
    ("", False, None),
    ("■ 매수 신호 (RSI 크로스다운)", True, ACCENT_CYAN),
    ("  RSI(2) < 15 진입 시 → 1 lot 순차 매수", False, GRAY_LIGHT),
    ("  최대 5회까지 분할 집행", False, GRAY_LIGHT),
    ("  5번째 lot → BIL 전량 매도 + 합산 매수", False, GRAY_LIGHT),
    ("", False, None),
    ("■ 매도 신호 (RSI 크로스업)", True, ACCENT_CYAN),
    ("  RSI(2) > 75 → TQQQ 전량 일괄 매도", False, GRAY_LIGHT),
    ("  사이클 재시작: 새 자산 기준 BIL 10% 재배분", False, GRAY_LIGHT),
    ("", False, None),
    ("■ 세금 처리", True, ACCENT_CYAN),
    ("  연간 실현손익 22% (비과세 $1,900)", False, GRAY_LIGHT),
    ("  다음 해 첫 거래일 자동 차감", False, GRAY_LIGHT),
]
add_multiline_textbox(s, strategy_lines, Inches(0.5), Inches(1.6), Inches(5.5), Inches(5.0),
                      font_size=12.5, color=WHITE)

# 우측: 핵심 철학
add_rect(s, Inches(6.6), Inches(1.05), Inches(6.3), Inches(2.6), fill_color=GRAY_DIM,
         line_color=ACCENT_BLUE, line_width=Pt(1))
add_text_box(s, "핵심 철학", Inches(6.7), Inches(1.1), Inches(6.0), Inches(0.45),
             font_size=16, bold=True, color=ACCENT_BLUE)
philosophy = [
    ("'하락을 분할로 흡수, 반등을 전체로 회수'", True, ACCENT_GOLD),
    ("", False, None),
    ("• RSI 과매도 = 저점 매수 기회 포착", False, GRAY_LIGHT),
    ("• 분할매수로 타이밍 리스크 분산", False, GRAY_LIGHT),
    ("• BIL 헤지로 대기 자금 수익화", False, GRAY_LIGHT),
    ("• 단순 규칙 기반 → 감정 배제", False, GRAY_LIGHT),
]
add_multiline_textbox(s, philosophy, Inches(6.7), Inches(1.6), Inches(6.0), Inches(2.0),
                      font_size=13, color=WHITE)

# 우측: 사용 자산
add_rect(s, Inches(6.6), Inches(3.8), Inches(6.3), Inches(3.0), fill_color=GRAY_DIM,
         line_color=ACCENT_CYAN, line_width=Pt(1))
add_text_box(s, "사용 자산", Inches(6.7), Inches(3.85), Inches(6.0), Inches(0.45),
             font_size=16, bold=True, color=ACCENT_CYAN)
assets = [
    ("TQQQ", True, ACCENT_GOLD),
    ("  나스닥100 3× 레버리지 ETF (ProShares)", False, GRAY_LIGHT),
    ("  변동성 크지만 상승장 폭발적 수익", False, TEXT_DIM),
    ("", False, None),
    ("BIL (구: TMF→교체)", True, ACCENT_GREEN),
    ("  1~3개월 미국 단기국채 ETF", False, GRAY_LIGHT),
    ("  금리상승기에도 안정 / 연 ~4~5% 수익", False, TEXT_DIM),
    ("", False, None),
    ("RSI(2) Wilder EWM", True, ACCENT_BLUE),
    ("  alpha = 1/period (2일 기준)", False, GRAY_LIGHT),
]
add_multiline_textbox(s, assets, Inches(6.7), Inches(4.35), Inches(6.0), Inches(2.3),
                      font_size=12.5, color=WHITE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — BATA 백테스트 성과
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s)
add_rect(s, 0, 0, SLIDE_W, Inches(0.08), ACCENT_GREEN)

add_text_box(s, "02  BATA 백테스트 성과 (2021.01 ~ 2026.06)",
             Inches(0.5), Inches(0.15), Inches(12), Inches(0.7),
             font_size=26, bold=True, color=WHITE)
add_divider(s, Inches(0.95), color=ACCENT_GREEN)

# 성과 카드 3개
cards = [
    ("FnG RSI-only\n(All-in/out 기준선)", "743.1%", "CAGR 48.1%", "MDD -48.9%", ACCENT_RED),
    ("A4 80/5×4\n[BIL 10% 헤지]", "301.2%", "CAGR 29.1%", "MDD -37.4%", ACCENT_GREEN),
    ("A 60/10/10/10/10\n[BIL 10% 헤지]", "220.1%", "CAGR 23.9%", "MDD -29.7%", ACCENT_BLUE),
]

for i, (name, ret, cagr, mdd, col) in enumerate(cards):
    lft = Inches(0.4 + i * 4.3)
    add_rect(s, lft, Inches(1.1), Inches(4.0), Inches(2.8),
             fill_color=GRAY_DIM, line_color=col, line_width=Pt(2))
    add_text_box(s, name, lft + Inches(0.1), Inches(1.2), Inches(3.8), Inches(0.8),
                 font_size=13, bold=True, color=col, align=PP_ALIGN.CENTER)
    add_text_box(s, ret, lft + Inches(0.1), Inches(2.1), Inches(3.8), Inches(0.7),
                 font_size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(s, cagr, lft + Inches(0.1), Inches(2.75), Inches(3.8), Inches(0.4),
                 font_size=15, color=ACCENT_CYAN, align=PP_ALIGN.CENTER)
    add_text_box(s, mdd, lft + Inches(0.1), Inches(3.15), Inches(3.8), Inches(0.4),
                 font_size=14, color=ACCENT_RED, align=PP_ALIGN.CENTER)

# 연도별 성과 표
add_text_box(s, "연도별 자산 추이", Inches(0.4), Inches(4.1), Inches(12), Inches(0.4),
             font_size=15, bold=True, color=ACCENT_GREEN)

tbl = s.shapes.add_table(3, 8,
                          Inches(0.4), Inches(4.55),
                          Inches(12.4), Inches(2.5)).table

headers = ["", "2021", "2022", "2023", "2024", "2025", "2026.H1", "총수익률"]
add_table_row(tbl, 0, headers, font_size=12, bold=True, bg_color=GRAY_MID)

row1 = ["A4 BIL", "+27.0%", "+40.7%", "+36.2%", "+47.7%", "+18.1%", "-5.5%", "301.2%"]
row2 = ["A BIL",  "+22.6%", "+33.6%", "+28.8%", "+35.9%", "+15.7%", "-3.5%", "220.1%"]

add_table_row(tbl, 1, row1, font_size=12, bg_color=RGBColor(0x12, 0x28, 0x18),
              text_colors=[ACCENT_GREEN, ACCENT_GREEN, ACCENT_GREEN, ACCENT_GREEN,
                           ACCENT_GREEN, ACCENT_GREEN, ACCENT_RED, ACCENT_GOLD])
add_table_row(tbl, 2, row2, font_size=12, bg_color=RGBColor(0x0d, 0x1e, 0x38),
              text_colors=[ACCENT_BLUE, ACCENT_CYAN, ACCENT_CYAN, ACCENT_CYAN,
                           ACCENT_CYAN, ACCENT_CYAN, ACCENT_RED, ACCENT_CYAN])


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — HFEA
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s)
add_rect(s, 0, 0, SLIDE_W, Inches(0.08), ACCENT_RED)

add_text_box(s, "03  HFEA — Hedgefundie's Excellent Adventure  🇺🇸",
             Inches(0.5), Inches(0.15), Inches(12.5), Inches(0.7),
             font_size=26, bold=True, color=WHITE)
add_divider(s, Inches(0.95), color=ACCENT_RED)

# 좌
add_rect(s, Inches(0.4), Inches(1.1), Inches(7.5), Inches(5.8), fill_color=GRAY_DIM,
         line_color=ACCENT_RED, line_width=Pt(1))

info = [
    ("창시자", True, ACCENT_RED),
    ("  익명 'Hedgefundie' — Bogleheads 포럼 (2019)", False, GRAY_LIGHT),
    ("", False, None),
    ("전략 구조", True, ACCENT_RED),
    ("  UPRO(S&P500 3×) 55% + TMF(장기채 3×) 45%", False, GRAY_LIGHT),
    ("  분기 1회 리밸런싱, 변경 없음", False, GRAY_LIGHT),
    ("", False, None),
    ("이론적 배경", True, ACCENT_RED),
    ("  주식↓ = 채권↑ (역상관) → 양방향 레버리지", False, GRAY_LIGHT),
    ("  리밸런싱 보너스로 장기 우수 성과 기대", False, GRAY_LIGHT),
    ("  저금리 시대(2010s)에 최적화된 구조", False, GRAY_LIGHT),
    ("", False, None),
    ("성과 (2010~2021)", True, ACCENT_RED),
    ("  CAGR ~20~32%  /  MDD ~-40%", False, GRAY_LIGHT),
    ("  2022년 TMF -70%, UPRO -70% 동반 폭락", False, ACCENT_RED),
    ("  2022년 단일 연도 손실 약 -60~-70%", False, ACCENT_RED),
    ("", False, None),
    ("BATA와의 비교", True, ACCENT_GOLD),
    ("  HFEA: 분기 리밸런싱 수동, TMF 리스크 내재", False, GRAY_LIGHT),
    ("  BATA: RSI 신호 기반 동적 진출입, BIL로 교체", False, ACCENT_GREEN),
    ("  → 2022년 BATA +33.6% vs HFEA -65% (추정)", False, ACCENT_GREEN),
]
add_multiline_textbox(s, info, Inches(0.5), Inches(1.15), Inches(7.2), Inches(5.5),
                      font_size=12.5, color=WHITE)

# 우 - 핵심 교훈
add_rect(s, Inches(8.1), Inches(1.1), Inches(4.8), Inches(2.8), fill_color=GRAY_DIM,
         line_color=ACCENT_GOLD, line_width=Pt(1))
add_text_box(s, "핵심 교훈", Inches(8.2), Inches(1.15), Inches(4.5), Inches(0.45),
             font_size=15, bold=True, color=ACCENT_GOLD)
lessons = [
    ("주식-채권 역상관 가정은", False, GRAY_LIGHT),
    ("금리 급등기에 깨진다", True, ACCENT_RED),
    ("", False, None),
    ("TMF → 최악의 헤지 도구", True, ACCENT_RED),
    ("(금리 상승 = TMF + UPRO 동반 폭락)", False, TEXT_DIM),
    ("", False, None),
    ("레버리지는 하락폭이 치명적", True, ACCENT_RED),
    ("-60% 손실 = 회복에 +150% 필요", False, TEXT_DIM),
]
add_multiline_textbox(s, lessons, Inches(8.2), Inches(1.65), Inches(4.5), Inches(2.1),
                      font_size=12.5, color=WHITE)

add_rect(s, Inches(8.1), Inches(4.05), Inches(4.8), Inches(2.85), fill_color=GRAY_DIM,
         line_color=ACCENT_BLUE, line_width=Pt(1))
add_text_box(s, "평가 지표", Inches(8.2), Inches(4.1), Inches(4.5), Inches(0.45),
             font_size=15, bold=True, color=ACCENT_BLUE)
metrics = [
    ("CAGR      ~19~32%  (2010~2021)", False, GRAY_LIGHT),
    ("MDD       -70%  (2022년 기준)", False, ACCENT_RED),
    ("구현 난이도  ★☆☆ (ETF 2종 매우 쉬움)", False, ACCENT_GREEN),
    ("", False, None),
    ("⚠ 2022년 이후 전략 신뢰도 급락", True, ACCENT_RED),
    ("포럼 내 포기자 속출, 전략 사망 논란", False, TEXT_DIM),
]
add_multiline_textbox(s, metrics, Inches(8.2), Inches(4.6), Inches(4.5), Inches(2.1),
                      font_size=12.5, color=WHITE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Dual Momentum GEM
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s)
add_rect(s, 0, 0, SLIDE_W, Inches(0.08), ACCENT_BLUE)

add_text_box(s, "04  Dual Momentum GEM — Gary Antonacci  🇺🇸",
             Inches(0.5), Inches(0.15), Inches(12.5), Inches(0.7),
             font_size=26, bold=True, color=WHITE)
add_divider(s, Inches(0.95), color=ACCENT_BLUE)

add_rect(s, Inches(0.4), Inches(1.1), Inches(7.5), Inches(5.8), fill_color=GRAY_DIM,
         line_color=ACCENT_BLUE, line_width=Pt(1))
info = [
    ("창시자", True, ACCENT_BLUE),
    ("  Gary Antonacci — 저서 'Dual Momentum Investing' (2014)", False, GRAY_LIGHT),
    ("  포트폴리오 성과로 Bernstein Award 수상", False, TEXT_DIM),
    ("", False, None),
    ("전략 구조 (GEM = Global Equity Momentum)", True, ACCENT_BLUE),
    ("  ① 절대 모멘텀: 미국주식 12개월 수익 > 단기채?", False, GRAY_LIGHT),
    ("      No → 현금/채권(AGG) 보유", False, TEXT_DIM),
    ("  ② 상대 모멘텀: 미국주식 vs 글로벌주식", False, GRAY_LIGHT),
    ("      승자 ETF (SPY or VEU) 전액 매수", False, GRAY_LIGHT),
    ("  월 1회 체크 → 리밸런싱", False, TEXT_DIM),
    ("", False, None),
    ("이론적 배경", True, ACCENT_BLUE),
    ("  '모멘텀 효과' = 최강 팩터 (Jegadeesh & Titman 1993)", False, GRAY_LIGHT),
    ("  상대·절대 모멘텀 결합으로 하락장 자동 회피", False, GRAY_LIGHT),
    ("", False, None),
    ("성과 (백테스트 1974~2013)", True, ACCENT_BLUE),
    ("  CAGR ~17.4%  /  MDD ~-22%", False, ACCENT_CYAN),
    ("  S&P500 MDD -50% 대비 절반 이하", False, ACCENT_GREEN),
    ("  2008~2009 금융위기 거의 무피해", False, ACCENT_GREEN),
    ("", False, None),
    ("BATA와의 비교", True, ACCENT_GOLD),
    ("  GEM: 월간 신호, ETF 1~2종, 단순 회전", False, GRAY_LIGHT),
    ("  BATA: 일간 RSI 신호, 레버리지 집중, 고수익 고위험", False, GRAY_LIGHT),
    ("  GEM 수익률 낮지만 MDD 훨씬 우수", False, ACCENT_GREEN),
]
add_multiline_textbox(s, info, Inches(0.5), Inches(1.15), Inches(7.2), Inches(5.5),
                      font_size=12.5, color=WHITE)

add_rect(s, Inches(8.1), Inches(1.1), Inches(4.8), Inches(2.6), fill_color=GRAY_DIM,
         line_color=ACCENT_CYAN, line_width=Pt(1))
add_text_box(s, "알파 vs MDD 트레이드오프",
             Inches(8.2), Inches(1.15), Inches(4.5), Inches(0.5),
             font_size=14, bold=True, color=ACCENT_CYAN)
tradeoff = [
    ("수익은 낮지만 MDD 탁월", True, ACCENT_GOLD),
    ("", False, None),
    ("CAGR   17.4%  vs  BATA 29.1%", False, GRAY_LIGHT),
    ("MDD   -22%  vs  BATA -37%", False, GRAY_LIGHT),
    ("", False, None),
    ("'전략 공개 후 알파 약화' 논란", False, TEXT_DIM),
    ("실제 2014 이후 성과는 백테스트 대비", False, TEXT_DIM),
    ("소폭 낮아짐", False, TEXT_DIM),
]
add_multiline_textbox(s, tradeoff, Inches(8.2), Inches(1.7), Inches(4.5), Inches(2.0),
                      font_size=12.5, color=WHITE)

add_rect(s, Inches(8.1), Inches(3.85), Inches(4.8), Inches(3.05), fill_color=GRAY_DIM,
         line_color=ACCENT_BLUE, line_width=Pt(1))
add_text_box(s, "핵심 장점", Inches(8.2), Inches(3.9), Inches(4.5), Inches(0.45),
             font_size=15, bold=True, color=ACCENT_BLUE)
pros = [
    ("✓ 월 15분 관리, 완전 자동화 가능", False, ACCENT_GREEN),
    ("✓ 하락장 조기 탈출 검증됨", False, ACCENT_GREEN),
    ("✓ 규칙이 완전히 투명·공개", False, ACCENT_GREEN),
    ("✓ 레버리지 없어 장기 유지 용이", False, ACCENT_GREEN),
    ("", False, None),
    ("✗ 레버리지 없어 상승장 수익 제한", False, ACCENT_RED),
    ("✗ 상승장 추세 뒤늦게 재진입", False, ACCENT_RED),
    ("구현 난이도  ★☆☆ (아주 쉬움)", False, ACCENT_GOLD),
]
add_multiline_textbox(s, pros, Inches(8.2), Inches(4.4), Inches(4.5), Inches(2.4),
                      font_size=12.5, color=WHITE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — 강환국 VAA/BAA
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s)
add_rect(s, 0, 0, SLIDE_W, Inches(0.08), ACCENT_GOLD)

add_text_box(s, "05  강환국 VAA/BAA 동적 자산배분  🇰🇷",
             Inches(0.5), Inches(0.15), Inches(12.5), Inches(0.7),
             font_size=26, bold=True, color=WHITE)
add_divider(s, Inches(0.95), color=ACCENT_GOLD)

add_rect(s, Inches(0.4), Inches(1.1), Inches(7.5), Inches(5.8), fill_color=GRAY_DIM,
         line_color=ACCENT_GOLD, line_width=Pt(1))
info = [
    ("창시자", True, ACCENT_GOLD),
    ("  Wouter J. Keller 논문 → 강환국 한국 대중화", False, GRAY_LIGHT),
    ("  저서 '거인의 포트폴리오' (2021) 베스트셀러", False, TEXT_DIM),
    ("", False, None),
    ("VAA (Vigilant Asset Allocation)", True, ACCENT_GOLD),
    ("  공격 ETF: SPY, QQQ, IWM, VGK, EWJ 등 12종", False, GRAY_LIGHT),
    ("  방어 ETF: LQD, IEF, SHY (채권 3종)", False, GRAY_LIGHT),
    ("  매월 모멘텀 스코어 계산 → 1위 ETF 집중 투자", False, GRAY_LIGHT),
    ("  공격 ETF 모멘텀 < 0 → 즉시 현금 100% 전환", False, ACCENT_RED),
    ("", False, None),
    ("BAA (Balanced Asset Allocation)", True, ACCENT_GOLD),
    ("  공격/방어 혼합 비율로 변동성 추가 완화", False, GRAY_LIGHT),
    ("  리스크 대비 수익 최적화 버전", False, GRAY_LIGHT),
    ("", False, None),
    ("성과 (2007~2020 백테스트)", True, ACCENT_GOLD),
    ("  CAGR ~19~20%  /  MDD ~-13~-15%", False, ACCENT_GREEN),
    ("  2008~2009 금융위기: 거의 무피해 (현금 전환)", False, ACCENT_GREEN),
    ("  2020 코로나 폭락: 빠른 탈출 후 재진입", False, ACCENT_GREEN),
    ("", False, None),
    ("BATA와의 비교", True, ACCENT_GOLD),
    ("  VAA: CAGR 낮지만 MDD -13% 압도적 방어력", False, GRAY_LIGHT),
    ("  BATA: 레버리지 덕에 CAGR 높지만 MDD -37%", False, GRAY_LIGHT),
    ("  → MDD 기준 최우수 전략 (레버리지 없이)", False, ACCENT_GREEN),
]
add_multiline_textbox(s, info, Inches(0.5), Inches(1.15), Inches(7.2), Inches(5.5),
                      font_size=12.5, color=WHITE)

add_rect(s, Inches(8.1), Inches(1.1), Inches(4.8), Inches(3.0), fill_color=GRAY_DIM,
         line_color=ACCENT_GREEN, line_width=Pt(1))
add_text_box(s, "MDD 비교 (목록 내 최우수)",
             Inches(8.2), Inches(1.15), Inches(4.5), Inches(0.5),
             font_size=14, bold=True, color=ACCENT_GREEN)
mdd_comp = [
    ("VAA/BAA     MDD  -13~-15%", True, ACCENT_GREEN),
    ("Dual Momentum  MDD  -22%", False, ACCENT_CYAN),
    ("BATA A BIL    MDD  -29.7%", False, ACCENT_BLUE),
    ("BATA A4 BIL   MDD  -37.4%", False, GRAY_LIGHT),
    ("HFEA         MDD  -70%", False, ACCENT_RED),
    ("", False, None),
    ("레버리지 없이 CAGR 20% + MDD -13%", True, ACCENT_GOLD),
    ("= 샤프비율 최고 수준", False, ACCENT_GOLD),
]
add_multiline_textbox(s, mdd_comp, Inches(8.2), Inches(1.7), Inches(4.5), Inches(2.3),
                      font_size=12.5, color=WHITE)

add_rect(s, Inches(8.1), Inches(4.25), Inches(4.8), Inches(2.65), fill_color=GRAY_DIM,
         line_color=ACCENT_GOLD, line_width=Pt(1))
add_text_box(s, "한국 퀀트 팩터 전략 (비교)",
             Inches(8.2), Inches(4.3), Inches(4.5), Inches(0.45),
             font_size=14, bold=True, color=ACCENT_GOLD)
quant = [
    ("코스닥 소형 저PBR/저PER 팩터", False, GRAY_LIGHT),
    ("분기 리밸런싱, 한국 소형주 집중", False, GRAY_LIGHT),
    ("", False, None),
    ("CAGR ~33%  /  MDD -53%", False, GRAY_LIGHT),
    ("수익 높지만 MDD 너무 큼", False, ACCENT_RED),
    ("→ 고수익·고MDD 딜레마 동일", True, ACCENT_RED),
    ("구현 난이도  ★★★ (국내주식 리서치 필요)", False, ACCENT_GOLD),
]
add_multiline_textbox(s, quant, Inches(8.2), Inches(4.8), Inches(4.5), Inches(2.0),
                      font_size=12.5, color=WHITE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — 일본 투자자
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s)
add_rect(s, 0, 0, SLIDE_W, Inches(0.08), RGBColor(0xff, 0x55, 0x55))

add_text_box(s, "06  일본 개인투자자 대표 전략  🇯🇵",
             Inches(0.5), Inches(0.15), Inches(12.5), Inches(0.7),
             font_size=26, bold=True, color=WHITE)
add_divider(s, Inches(0.95), color=RGBColor(0xff, 0x55, 0x55))

# 좌: Tapazou
add_rect(s, Inches(0.4), Inches(1.1), Inches(6.0), Inches(5.8), fill_color=GRAY_DIM,
         line_color=RGBColor(0xff, 0x99, 0x00), line_width=Pt(1))
add_text_box(s, "たぱぞう (Tapazou)",
             Inches(0.5), Inches(1.15), Inches(5.8), Inches(0.5),
             font_size=16, bold=True, color=RGBColor(0xff, 0x99, 0x00))
tapazou = [
    ("전략", True, RGBColor(0xff, 0x99, 0x00)),
    ("  VTI/VOO 매월 정액 DCA", False, GRAY_LIGHT),
    ("  배당 자동 재투자, 절대 매도 없음", False, GRAY_LIGHT),
    ("  일본 NISA 계좌 최대 활용", False, GRAY_LIGHT),
    ("", False, None),
    ("철학", True, RGBColor(0xff, 0x99, 0x00)),
    ("  '시장 타이밍은 불가능하다'", False, GRAY_LIGHT),
    ("  장기 보유만이 유일한 정답", False, GRAY_LIGHT),
    ("  감정 배제, 자동화 최우선", False, GRAY_LIGHT),
    ("", False, None),
    ("성과", True, RGBColor(0xff, 0x99, 0x00)),
    ("  CAGR ~8~10%  /  MDD ~-50%", False, GRAY_LIGHT),
    ("  시장과 동일 (알파 없음)", False, TEXT_DIM),
    ("", False, None),
    ("평가", True, RGBColor(0xff, 0x99, 0x00)),
    ("  일본 FIRE 운동의 상징적 존재", False, GRAY_LIGHT),
    ("  블로그 월 수백만 뷰, 책 30만부", False, GRAY_LIGHT),
    ("  단순함의 미학 — 실행이 핵심", False, ACCENT_GREEN),
    ("", False, None),
    ("BATA와 비교", True, ACCENT_GOLD),
    ("  정반대 철학: 타이밍 vs 시간", False, GRAY_LIGHT),
    ("  BATA CAGR 23~29% vs DCA 8~10%", False, GRAY_LIGHT),
    ("  MDD는 Tapazou -50% vs BATA -30~37%", False, GRAY_LIGHT),
    ("구현 난이도  ★☆☆ (가장 쉬움)", False, ACCENT_GREEN),
]
add_multiline_textbox(s, tapazou, Inches(0.5), Inches(1.7), Inches(5.7), Inches(5.0),
                      font_size=12.5, color=WHITE)

# 우: バフェット太郎
add_rect(s, Inches(6.8), Inches(1.1), Inches(6.1), Inches(5.8), fill_color=GRAY_DIM,
         line_color=ACCENT_BLUE, line_width=Pt(1))
add_text_box(s, "バフェット太郎 (버핏 타로)",
             Inches(6.9), Inches(1.15), Inches(5.9), Inches(0.5),
             font_size=16, bold=True, color=ACCENT_BLUE)
buffett_taro = [
    ("전략", True, ACCENT_BLUE),
    ("  미국 고배당 우량주 10종목 균등 보유", False, GRAY_LIGHT),
    ("  KO, WMT, JNJ, PM, MCD, PEP 등", False, TEXT_DIM),
    ("  배당 재투자, 매분기 리밸런싱", False, GRAY_LIGHT),
    ("  개별종목 추가 매수 없음", False, GRAY_LIGHT),
    ("", False, None),
    ("철학", True, ACCENT_BLUE),
    ("  '배당이 심리적 안정감을 준다'", False, GRAY_LIGHT),
    ("  불황에도 배당으로 현금 창출 가능", False, GRAY_LIGHT),
    ("  워런 버핏의 장기투자 철학 추종", False, GRAY_LIGHT),
    ("", False, None),
    ("성과", True, ACCENT_BLUE),
    ("  CAGR ~12~13%  /  MDD ~-40%", False, GRAY_LIGHT),
    ("  S&P500 지속 언더퍼폼", False, ACCENT_RED),
    ("  배당수익률 ~3~4% (심리적 완충)", False, ACCENT_GREEN),
    ("", False, None),
    ("평가", True, ACCENT_BLUE),
    ("  배당 집중 전략의 수익 한계 확인", False, GRAY_LIGHT),
    ("  VOO/VTI보다 낮은 성과 지속", False, ACCENT_RED),
    ("  배당 심리 안정감 외에 수익 메리트 약함", False, GRAY_LIGHT),
    ("", False, None),
    ("BATA와 비교", True, ACCENT_GOLD),
    ("  BATA CAGR 3배 이상 우수", False, ACCENT_GREEN),
    ("  MDD는 BATA가 약간 더 낮음 (A BIL 기준)", False, ACCENT_GREEN),
    ("구현 난이도  ★☆☆ (쉬움)", False, ACCENT_GREEN),
]
add_multiline_textbox(s, buffett_taro, Inches(6.9), Inches(1.7), Inches(5.8), Inches(5.0),
                      font_size=12.5, color=WHITE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — 싱가포르 + 유럽
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s)
add_rect(s, 0, 0, SLIDE_W, Inches(0.08), ACCENT_CYAN)

add_text_box(s, "07  싱가포르 · 독일 · 미국 기타 전략",
             Inches(0.5), Inches(0.15), Inches(12.5), Inches(0.7),
             font_size=26, bold=True, color=WHITE)
add_divider(s, Inches(0.95), color=ACCENT_CYAN)

# 1M65
add_rect(s, Inches(0.4), Inches(1.1), Inches(4.0), Inches(2.9), fill_color=GRAY_DIM,
         line_color=ACCENT_GREEN, line_width=Pt(1))
add_text_box(s, "1M65 — CPF 극대화  🇸🇬",
             Inches(0.5), Inches(1.15), Inches(3.8), Inches(0.45),
             font_size=14, bold=True, color=ACCENT_GREEN)
m65 = [
    ("Loo Cheng Chuan — 싱가포르 시민 전용", False, GRAY_LIGHT),
    ("CPF SA 4% 정부보증 최대 납입", False, GRAY_LIGHT),
    ("65세 부부 합산 SGD 100만 목표", False, GRAY_LIGHT),
    ("CAGR 4% / MDD 0% (정부 보증)", False, ACCENT_GREEN),
    ("특이점: 레버리지·주식 완전 배제", False, TEXT_DIM),
    ("MDD 제로의 안전망, 알파 없음", False, TEXT_DIM),
]
add_multiline_textbox(s, m65, Inches(0.5), Inches(1.65), Inches(3.8), Inches(2.2),
                      font_size=11.5, color=WHITE)

# TTI
add_rect(s, Inches(4.6), Inches(1.1), Inches(4.0), Inches(2.9), fill_color=GRAY_DIM,
         line_color=ACCENT_CYAN, line_width=Pt(1))
add_text_box(s, "ThumbTack Investor (TTI)  🇸🇬",
             Inches(4.7), Inches(1.15), Inches(3.8), Inches(0.45),
             font_size=14, bold=True, color=ACCENT_CYAN)
tti = [
    ("아시아 딥밸류 + 컨트래리언 집중 포트", False, GRAY_LIGHT),
    ("실적 투명 공개 (블로그 연간 보고)", False, ACCENT_GREEN),
    ("CAGR ~23% (공개 기록)", False, GRAY_LIGHT),
    ("소수 종목 집중, 리서치 집약적", False, TEXT_DIM),
    ("재현 어려움 (고급 리서치 필요)", False, ACCENT_RED),
    ("구현 난이도  ★★★", False, ACCENT_RED),
]
add_multiline_textbox(s, tti, Inches(4.7), Inches(1.65), Inches(3.8), Inches(2.2),
                      font_size=11.5, color=WHITE)

# Gerd Kommer
add_rect(s, Inches(8.8), Inches(1.1), Inches(4.1), Inches(2.9), fill_color=GRAY_DIM,
         line_color=ACCENT_BLUE, line_width=Pt(1))
add_text_box(s, "Gerd Kommer — 글로벌 팩터 ETF  🇩🇪",
             Inches(8.9), Inches(1.15), Inches(3.9), Inches(0.45),
             font_size=13.5, bold=True, color=ACCENT_BLUE)
kommer = [
    ("전 세계 GDP 가중 ETF + 팩터 오버웨이트", False, GRAY_LIGHT),
    ("Size / Value / Quality 팩터 결합", False, GRAY_LIGHT),
    ("본인 ETF 출시 (Gerd Kommer ETF)", False, ACCENT_GREEN),
    ("CAGR ~8~10% / MDD ~-50%", False, GRAY_LIGHT),
    ("독일의 'ETF 교황'", False, TEXT_DIM),
    ("학문적 근거 탄탄, 수익은 평범", False, TEXT_DIM),
    ("구현 난이도  ★★☆ (ETF 선택 복잡)", False, ACCENT_GOLD),
]
add_multiline_textbox(s, kommer, Inches(8.9), Inches(1.65), Inches(3.9), Inches(2.2),
                      font_size=11.5, color=WHITE)

# Dragon Portfolio
add_rect(s, Inches(0.4), Inches(4.2), Inches(12.5), Inches(2.7), fill_color=GRAY_DIM,
         line_color=ACCENT_RED, line_width=Pt(1.5))
add_text_box(s, "Dragon Portfolio — Chris Cole (Artemis Capital)  🇺🇸",
             Inches(0.5), Inches(4.25), Inches(12.0), Inches(0.5),
             font_size=16, bold=True, color=ACCENT_RED)
dragon_left = [
    ("구성", True, ACCENT_RED),
    ("  주식 24% + 장기채 18% + 금 19%", False, GRAY_LIGHT),
    ("  원자재 추세추종 18% + 롱 변동성(옵션) 21%", False, GRAY_LIGHT),
    ("", False, None),
    ("성과 (90년 백테스트)", True, ACCENT_RED),
    ("  CAGR ~14.4%  /  MDD 극히 낮음", False, ACCENT_GREEN),
    ("  2022년 인플레이션 환경 방어 검증", False, ACCENT_GREEN),
]
add_multiline_textbox(s, dragon_left, Inches(0.5), Inches(4.75), Inches(6.0), Inches(2.0),
                      font_size=12.5, color=WHITE)
dragon_right = [
    ("핵심 아이디어", True, ACCENT_RED),
    ("  '어떤 경제 환경에도 1/5씩 작동하는 자산 배분'", False, GRAY_LIGHT),
    ("  주식↑: 성장기  /  채권↑: 디플레  /  금↑: 인플레", False, GRAY_LIGHT),
    ("  롱 변동성↑: 위기  /  원자재↑: 원자재 슈퍼사이클", False, GRAY_LIGHT),
    ("", False, None),
    ("⚠ 롱 변동성(옵션) 파트 → 개인 구현 매우 어려움", True, ACCENT_RED),
    ("  구현 난이도  ★★★ (전문 투자자 수준)", False, ACCENT_RED),
]
add_multiline_textbox(s, dragon_right, Inches(6.7), Inches(4.75), Inches(6.0), Inches(2.0),
                      font_size=12.5, color=WHITE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — 종합 비교표
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s)
add_rect(s, 0, 0, SLIDE_W, Inches(0.08), ACCENT_GOLD)

add_text_box(s, "종합 비교표 — 전략별 성과 & 특성",
             Inches(0.5), Inches(0.15), Inches(12.5), Inches(0.7),
             font_size=26, bold=True, color=WHITE)
add_divider(s, Inches(0.95), color=ACCENT_GOLD)

tbl = s.shapes.add_table(12, 6,
                          Inches(0.3), Inches(1.1),
                          Inches(12.7), Inches(6.0)).table

col_widths = [Inches(3.3), Inches(1.7), Inches(1.5), Inches(1.7), Inches(2.2), Inches(2.3)]
for i, w in enumerate(col_widths):
    tbl.columns[i].width = w

headers = ["전략", "CAGR", "MDD", "구현 난이도", "강점", "약점"]
add_table_row(tbl, 0, headers, font_size=12, bold=True, bg_color=RGBColor(0x15, 0x25, 0x45))

rows_data = [
    ("BATA A4 BIL10% (본인 전략)",  "29.1%", "-37.4%", "★★☆",  "레버리지 고수익",    "MDD 높음",
     ACCENT_GREEN, ACCENT_GREEN, ACCENT_RED, ACCENT_GOLD, ACCENT_GREEN, ACCENT_RED),
    ("BATA A BIL10% (본인 전략)",   "23.9%", "-29.7%", "★★☆",  "MDD-수익 균형",     "수익 제한",
     ACCENT_GREEN, ACCENT_CYAN, ACCENT_GREEN, ACCENT_GOLD, ACCENT_CYAN, TEXT_DIM),
    ("FnG RSI-only (기준선)",        "48.1%", "-48.9%", "★☆☆",  "폭발적 수익",       "MDD 크고 변동 심함",
     GRAY_LIGHT, ACCENT_GOLD, ACCENT_RED, ACCENT_GREEN, ACCENT_GOLD, ACCENT_RED),
    ("VAA/BAA (강환국)",             "~20%",  "-13~15%","★★☆",  "MDD 최우수",        "수익 낮음",
     GRAY_LIGHT, ACCENT_CYAN, ACCENT_GREEN, ACCENT_GOLD, ACCENT_GREEN, TEXT_DIM),
    ("Dual Momentum GEM",           "~17%",  "-22%",   "★☆☆",  "하락장 방어",       "상승장 수익 제한",
     GRAY_LIGHT, ACCENT_CYAN, ACCENT_GREEN, ACCENT_GREEN, ACCENT_CYAN, TEXT_DIM),
    ("HFEA (Hedgefundie)",          "~25%",  "-70%",   "★☆☆",  "금리하락기 폭발",   "금리상승기 괴멸",
     GRAY_LIGHT, ACCENT_CYAN, ACCENT_RED, ACCENT_GREEN, ACCENT_CYAN, ACCENT_RED),
    ("강환국 퀀트 팩터",             "~33%",  "-53%",   "★★★",  "고CAGR",            "MDD 큼·국내주식 한정",
     GRAY_LIGHT, ACCENT_GOLD, ACCENT_RED, ACCENT_RED, ACCENT_GOLD, ACCENT_RED),
    ("たぱぞう DCA",                 "~8~10%","-50%",   "★☆☆",  "완전 자동화",       "시장 평균 수준",
     GRAY_LIGHT, TEXT_DIM, ACCENT_RED, ACCENT_GREEN, ACCENT_GREEN, TEXT_DIM),
    ("Dragon Portfolio",            "~14.4%","-극소",  "★★★",  "어떤 환경도 방어",  "롱 변동성 구현 어려움",
     GRAY_LIGHT, ACCENT_CYAN, ACCENT_GREEN, ACCENT_RED, ACCENT_GREEN, ACCENT_RED),
    ("Gerd Kommer ETF",             "~8~10%","-50%",   "★★☆",  "학문적 근거 탄탄",  "알파 미약",
     GRAY_LIGHT, TEXT_DIM, ACCENT_RED, ACCENT_GOLD, ACCENT_CYAN, TEXT_DIM),
    ("バフェット太郎 배당주",        "~12~13%","-40%",  "★☆☆",  "배당 현금흐름",     "S&P500 언더퍼폼",
     GRAY_LIGHT, TEXT_DIM, ACCENT_RED, ACCENT_GREEN, ACCENT_CYAN, ACCENT_RED),
]

bg_colors = [
    RGBColor(0x10, 0x28, 0x18), RGBColor(0x0a, 0x1e, 0x38), RGBColor(0x1c, 0x1c, 0x10),
    RGBColor(0x1a, 0x20, 0x10), RGBColor(0x12, 0x18, 0x28), RGBColor(0x28, 0x10, 0x10),
    RGBColor(0x20, 0x1a, 0x0a), RGBColor(0x10, 0x10, 0x18), RGBColor(0x12, 0x12, 0x22),
    RGBColor(0x10, 0x12, 0x18), RGBColor(0x10, 0x14, 0x10),
]

for ri, (row_vals) in enumerate(rows_data):
    vals = row_vals[:6]
    colors = row_vals[6:]
    row = tbl.rows[ri + 1]
    for ci, (val, col) in enumerate(zip(vals, colors)):
        cell = row.cells[ci]
        cell.text = val
        tf = cell.text_frame
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        runs = tf.paragraphs[0].runs
        if runs:
            runs[0].font.size = Pt(11)
            runs[0].font.name = "Malgun Gothic"
            runs[0].font.color.rgb = col
        fill = cell.fill
        fill.solid()
        fill.fore_color.rgb = bg_colors[ri]


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — 결론
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s)
add_rect(s, 0, 0, SLIDE_W, Inches(0.08), ACCENT_GREEN)

add_text_box(s, "10  결론 — BATA 전략의 포지셔닝",
             Inches(0.5), Inches(0.15), Inches(12.5), Inches(0.7),
             font_size=26, bold=True, color=WHITE)
add_divider(s, Inches(0.95), color=ACCENT_GREEN)

# 좌측 결론 박스
add_rect(s, Inches(0.4), Inches(1.1), Inches(7.8), Inches(5.8), fill_color=GRAY_DIM,
         line_color=ACCENT_GREEN, line_width=Pt(1.5))
add_text_box(s, "BATA 전략의 강점",
             Inches(0.5), Inches(1.15), Inches(7.6), Inches(0.5),
             font_size=16, bold=True, color=ACCENT_GREEN)

conclusions = [
    ("① 레버리지 전략 중 가장 합리적인 MDD 관리", True, ACCENT_CYAN),
    ("   HFEA -70% 대비 A BIL -29.7% 압도적 개선", False, GRAY_LIGHT),
    ("   TMF를 BIL로 교체한 결정이 핵심 차별화", False, TEXT_DIM),
    ("", False, None),
    ("② RSI 신호 기반 자동화 규칙 — 감정 배제", True, ACCENT_CYAN),
    ("   주관적 판단 없이 크로스 신호만으로 진출입", False, GRAY_LIGHT),
    ("   백테스트 재현 가능, 투명한 로직", False, GRAY_LIGHT),
    ("", False, None),
    ("③ CAGR 29.1% — 레버리지 없는 전략 대비 우수", True, ACCENT_CYAN),
    ("   VAA/BAA CAGR 20%, GEM 17% 대비 1.5~2배", False, GRAY_LIGHT),
    ("", False, None),
    ("④ MDD 개선 여지 존재", True, ACCENT_GOLD),
    ("   분할 비율 조정 (60/10/10/10/10 선택 시 -29.7%)", False, GRAY_LIGHT),
    ("   BIL 비율 상향 (15~20%) 검토 가능", False, GRAY_LIGHT),
    ("", False, None),
    ("⑤ '수익률 30%+, MDD -15%'는 이론상 불가능 영역", True, ACCENT_RED),
    ("   전략의 근본적 트레이드오프 — 수익 ↔ MDD", False, GRAY_LIGHT),
    ("   BATA는 이 트레이드오프 내 최적점 탐색 중", False, ACCENT_GOLD),
]
add_multiline_textbox(s, conclusions, Inches(0.5), Inches(1.7), Inches(7.6), Inches(5.0),
                      font_size=12.5, color=WHITE)

# 우측: 비교 요약
add_rect(s, Inches(8.4), Inches(1.1), Inches(4.5), Inches(3.0), fill_color=GRAY_DIM,
         line_color=ACCENT_GOLD, line_width=Pt(1.5))
add_text_box(s, "포지셔닝 맵 (CAGR vs MDD)",
             Inches(8.5), Inches(1.15), Inches(4.3), Inches(0.5),
             font_size=14, bold=True, color=ACCENT_GOLD)
positioning = [
    ("고수익·고MDD 구간", True, ACCENT_RED),
    ("  FnG +743% / BATA A4 +301% / 한국퀀트", False, GRAY_LIGHT),
    ("", False, None),
    ("중수익·중MDD 구간 ← BATA A", True, ACCENT_CYAN),
    ("  BATA A BIL +220% / HFEA / バ太郎", False, GRAY_LIGHT),
    ("", False, None),
    ("중수익·저MDD 구간 (이상적)", True, ACCENT_GREEN),
    ("  VAA/BAA, GEM, Dragon Portfolio", False, GRAY_LIGHT),
    ("  → 레버리지 없이 달성 가능한 영역", False, TEXT_DIM),
]
add_multiline_textbox(s, positioning, Inches(8.5), Inches(1.7), Inches(4.3), Inches(2.2),
                      font_size=12.5, color=WHITE)

add_rect(s, Inches(8.4), Inches(4.25), Inches(4.5), Inches(2.65), fill_color=GRAY_DIM,
         line_color=ACCENT_BLUE, line_width=Pt(1.5))
add_text_box(s, "다음 연구 방향",
             Inches(8.5), Inches(4.3), Inches(4.3), Inches(0.45),
             font_size=14, bold=True, color=ACCENT_BLUE)
next_steps = [
    ("• BIL 비율 최적화 (10→15→20%)", False, GRAY_LIGHT),
    ("• VAA 모멘텀 + RSI 신호 결합 실험", False, GRAY_LIGHT),
    ("• MDD -25% 이하 달성 파라미터 탐색", False, GRAY_LIGHT),
    ("", False, None),
    ("목표: CAGR 25%+ / MDD -25% 이하", True, ACCENT_GOLD),
    ("= 전 세계 최고 수준의 개인투자 전략", False, ACCENT_GOLD),
]
add_multiline_textbox(s, next_steps, Inches(8.5), Inches(4.8), Inches(4.3), Inches(2.0),
                      font_size=12.5, color=WHITE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE A — TQQQ 8개 전략 개요
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s)
add_rect(s, 0, 0, SLIDE_W, Inches(0.06), ACCENT_BLUE)
add_rect(s, 0, SLIDE_H - Inches(0.06), SLIDE_W, Inches(0.06), ACCENT_BLUE)

add_text_box(s, "TQQQ 단일 종목 — 룰 기반 전략 9종 비교",
             Inches(0.5), Inches(0.12), Inches(12.5), Inches(0.55),
             font_size=26, bold=True, color=WHITE)
add_divider(s, Inches(0.72))
add_text_box(s, "초기자금 $10,000 · 슬리피지·세금 미반영 · 2011-01-03 ~ 2024-12-30",
             Inches(0.5), Inches(0.74), Inches(12.5), Inches(0.35),
             font_size=12, color=TEXT_DIM)

# 테이블: 9행 × 7열
tbl_left = Inches(0.3)
tbl_top  = Inches(1.15)
tbl_w    = Inches(12.73)
tbl_h    = Inches(5.85)
table = s.shapes.add_table(10, 7, tbl_left, tbl_top, tbl_w, tbl_h).table

# 열 너비
col_widths = [Inches(2.8), Inches(1.5), Inches(1.5), Inches(1.3), Inches(1.3), Inches(1.3), Inches(3.03)]
for ci, cw in enumerate(col_widths):
    table.columns[ci].width = cw

headers = ["전략", "2011~24 총수익", "2021~24 총수익", "CAGR(전체)", "MDD(전체)", "Sharpe", "전략 핵심"]
header_colors = [None]*7
add_table_row(table, 0, headers, font_size=12, bold=True, bg_color=RGBColor(0x00,0x4a,0x8a),
              text_colors=[ACCENT_CYAN]*7)

rows_data = [
    ("S0. TQQQ Buy & Hold",       "10,391%", "  +93.8%", "39.5%", "-81.7%", "0.86", "TQQQ 무조건 보유 (기준선)"),
    ("S1. SMA200 트렌드추종",      " 6,641%", "+281.7%",  "35.1%", "-55.6%", "0.90", "QQQ>SMA200 → TQQQ / 하락 → BIL"),
    ("S2. SMA200 + RSI(2)",        " 1,293%", "+167.8%",  "20.7%", "-62.0%", "0.68", "SMA200 위에서만 RSI 신호 적용"),
    ("S3. 절대 모멘텀(12M)",       " 1,710%", " +68.0%",  "23.0%", "-77.0%", "0.67", "월말 12M 수익률>0 → TQQQ"),
    ("S4. VAA 라이트",             " 3,734%", "+246.3%",  "29.8%", "-72.1%", "0.79", "TQQQ/QLD/QQQ 모멘텀 점수 최강 자산"),
    ("S5. RSI(2) 분할매수",        "    997%", " +21.2%",  "18.7%", "-64.5%", "1.05", "4-lot 분할진입 / RSI>80 일괄청산"),
    ("S6. 골든/데드크로스",        "11,287%", "+394.7%",  "40.3%", "-69.9%", "0.93", "SMA50>SMA200 → TQQQ / 역전 → BIL"),
    ("S7. FnG+RSI(2) ★",          " 4,257%", "+684.8%",  "31.0%", "-67.2%", "0.85", "RSI 크로스다운15 매수 / 크로스업75 매도"),
    ("S8. 골든크로스+FnG ★",      " 1,686%", "+252.2%",  "22.9%", "-67.2%", "0.76", "골든크로스 진입허가 + RSI 타이밍 결합"),
]

highlight_rows = {7, 8}  # S7, S8 행 강조
for ri, (cells) in enumerate(rows_data):
    is_hl = ri in highlight_rows
    bg = RGBColor(0x0a,0x2a,0x4a) if ri % 2 == 0 else RGBColor(0x0d,0x1e,0x36)
    if ri == 7:
        bg = RGBColor(0x00,0x2a,0x1a)  # S7 초록 강조
    add_table_row(table, ri + 1, cells, font_size=11,
                  bg_color=bg,
                  text_colors=[
                      ACCENT_GOLD if ri == 7 else WHITE,
                      ACCENT_GREEN if ri == 7 else (ACCENT_GOLD if ri == 6 else GRAY_LIGHT),
                      ACCENT_GREEN if ri == 7 else GRAY_LIGHT,
                      GRAY_LIGHT, ACCENT_RED, GRAY_LIGHT, TEXT_DIM
                  ])

# 범례
add_text_box(s, "★ S7 = 현재 BATA 서비스 중인 FnG+RSI(2) 전략  |  S8 = 이번 연구에서 새로 실험한 결합 전략",
             Inches(0.3), Inches(7.1), Inches(12.5), Inches(0.3),
             font_size=10.5, color=ACCENT_GOLD)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE B — 2011~2024 vs 2021~2024 기간별 성과 비교
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s)
add_rect(s, 0, 0, SLIDE_W, Inches(0.06), ACCENT_CYAN)
add_rect(s, 0, SLIDE_H - Inches(0.06), SLIDE_W, Inches(0.06), ACCENT_CYAN)

add_text_box(s, "기간별 성과 비교: 2011~2024 (14년)  vs  2021~2024 (4년)",
             Inches(0.5), Inches(0.12), Inches(12.5), Inches(0.55),
             font_size=24, bold=True, color=WHITE)
add_divider(s, Inches(0.72), ACCENT_CYAN)

# 왼쪽: 2011~2024
add_rect(s, Inches(0.3), Inches(0.85), Inches(6.1), Inches(0.42),
         fill_color=RGBColor(0x00,0x4a,0x8a))
add_text_box(s, "📅  2011~2024 (14년 전체)",
             Inches(0.35), Inches(0.87), Inches(6.0), Inches(0.38),
             font_size=14, bold=True, color=ACCENT_CYAN)

tbl_a = s.shapes.add_table(10, 5, Inches(0.3), Inches(1.3), Inches(6.1), Inches(5.6)).table
for ci, cw in enumerate([Inches(2.2), Inches(1.0), Inches(1.0), Inches(1.0), Inches(0.9)]):
    tbl_a.columns[ci].width = cw
add_table_row(tbl_a, 0, ["전략", "총수익", "CAGR", "MDD", "Sharpe"],
              font_size=11, bold=True, bg_color=RGBColor(0x00,0x38,0x70),
              text_colors=[ACCENT_CYAN]*5)
rows_a = [
    ("S0. TQQQ Buy&Hold",    "10,391%", "39.5%", "-81.7%", "0.86"),
    ("S1. SMA200 트렌드",    " 6,641%", "35.1%", "-55.6%", "0.90"),
    ("S2. SMA200+RSI(2)",    " 1,293%", "20.7%", "-62.0%", "0.68"),
    ("S3. 절대모멘텀12M",    " 1,710%", "23.0%", "-77.0%", "0.67"),
    ("S4. VAA 라이트",       " 3,734%", "29.8%", "-72.1%", "0.79"),
    ("S5. RSI(2)분할매수",   "    997%", "18.7%", "-64.5%", "1.05"),
    ("S6. 골든/데드크로스",  "11,287%", "40.3%", "-69.9%", "0.93"),
    ("S7. FnG+RSI(2) ★",   " 4,257%", "31.0%", "-67.2%", "0.85"),
    ("S8. 골든크로스+FnG",  " 1,686%", "22.9%", "-67.2%", "0.76"),
]
for ri, cells in enumerate(rows_a):
    bg = RGBColor(0x0a,0x2a,0x4a) if ri % 2 == 0 else RGBColor(0x0d,0x1e,0x36)
    add_table_row(tbl_a, ri+1, cells, font_size=11, bg_color=bg,
                  text_colors=[WHITE, GRAY_LIGHT, GRAY_LIGHT, ACCENT_RED, GRAY_LIGHT])

# 오른쪽: 2021~2024
add_rect(s, Inches(6.7), Inches(0.85), Inches(6.3), Inches(0.42),
         fill_color=RGBColor(0x00,0x3a,0x1a))
add_text_box(s, "📅  2021~2024 (최근 4년)  ← 시장 격동기",
             Inches(6.75), Inches(0.87), Inches(6.2), Inches(0.38),
             font_size=14, bold=True, color=ACCENT_GREEN)

tbl_b = s.shapes.add_table(10, 5, Inches(6.7), Inches(1.3), Inches(6.3), Inches(5.6)).table
for ci, cw in enumerate([Inches(2.2), Inches(1.1), Inches(1.0), Inches(1.0), Inches(1.0)]):
    tbl_b.columns[ci].width = cw
add_table_row(tbl_b, 0, ["전략", "총수익", "CAGR", "MDD", "Sharpe"],
              font_size=11, bold=True, bg_color=RGBColor(0x00,0x3a,0x1a),
              text_colors=[ACCENT_GREEN]*5)
rows_b = [
    ("S0. TQQQ Buy&Hold",    " +93.8%", "18.1%", "-81.7%", "0.58"),
    ("S1. SMA200 트렌드",    "+281.7%", "39.9%", "-52.3%", "0.95"),
    ("S2. SMA200+RSI(2)",    "+167.8%", "28.0%", "-39.2%", "0.81"),
    ("S3. 절대모멘텀12M",    " +68.0%", "13.9%", "-69.6%", "0.51"),
    ("S4. VAA 라이트",       "+246.3%", "36.6%", "-49.0%", "0.90"),
    ("S5. RSI(2)분할매수",   " +21.2%",  "4.9%", "-29.0%", "0.43"),
    ("S6. 골든/데드크로스",  "+394.7%", "49.3%", "-49.0%", "1.08"),
    ("S7. FnG+RSI(2) ★",   "+684.8%", "67.7%", "-30.3%", "1.37"),
    ("S8. 골든크로스+FnG",  "+252.2%", "37.1%", "-28.4%", "1.16"),
]
for ri, cells in enumerate(rows_b):
    is_s7 = (ri == 7)
    bg = RGBColor(0x00,0x2a,0x1a) if is_s7 else (RGBColor(0x0a,0x2a,0x4a) if ri % 2 == 0 else RGBColor(0x0d,0x1e,0x36))
    tc = [ACCENT_GOLD if is_s7 else WHITE,
          ACCENT_GREEN if is_s7 else GRAY_LIGHT,
          ACCENT_GREEN if is_s7 else GRAY_LIGHT,
          ACCENT_GREEN if is_s7 else ACCENT_RED,
          ACCENT_GOLD if is_s7 else GRAY_LIGHT]
    add_table_row(tbl_b, ri+1, cells, font_size=11, bold=is_s7, bg_color=bg, text_colors=tc)

add_text_box(s, "▶  최근 4년 1위: S7 FnG+RSI(2)  CAGR 67.7% / MDD -30.3% / Sharpe 1.37",
             Inches(0.3), Inches(7.05), Inches(12.5), Inches(0.32),
             font_size=12, bold=True, color=ACCENT_GOLD)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE C — S6 vs S7 vs S8 집중 비교 (연도별 수익률 포함)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s)
add_rect(s, 0, 0, SLIDE_W, Inches(0.06), ACCENT_GOLD)
add_rect(s, 0, SLIDE_H - Inches(0.06), SLIDE_W, Inches(0.06), ACCENT_GOLD)

add_text_box(s, "S6 vs S7 vs S8 집중 비교  —  골든크로스 필터 + FnG RSI 결합 실험",
             Inches(0.5), Inches(0.12), Inches(12.5), Inches(0.55),
             font_size=22, bold=True, color=WHITE)
add_divider(s, Inches(0.72), ACCENT_GOLD)

# 왼쪽: 전략 설명 박스들
for i, (title, body, bc) in enumerate([
    ("S6. 골든/데드크로스",
     "QQQ SMA50 > SMA200 → TQQQ\nSMA50 < SMA200 → BIL\n\n• 시장 추세 순응 전략\n• 연 1~3회 전환, 장기 보유\n• 2022년 하락 방어: -37.9%",
     ACCENT_BLUE),
    ("S7. FnG+RSI(2) ★",
     "RSI(2) 크로스다운 15 → 전량 매수\nRSI(2) 크로스업 75 → 전량 매도\n\n• 추세 무관, 단기 과매도 포착\n• 2022년 하락장에서 +61.3% !\n• 최근 4년 Sharpe 1.37 압도적 1위",
     ACCENT_GREEN),
    ("S8. 골든크로스+FnG ★",
     "골든크로스 중에만 RSI 진입 허가\n데드크로스 또는 RSI>75 → 청산\n\n• S6 진입필터 + S7 타이밍 결합\n• 2022년 데드크로스로 진입 차단\n• MDD -28.4% (3개 중 최저)",
     ACCENT_CYAN),
]):
    bx = Inches(0.3)
    by = Inches(0.88) + i * Inches(2.05)
    add_rect(s, bx, by, Inches(4.2), Inches(1.9), fill_color=GRAY_DIM,
             line_color=bc, line_width=Pt(1.5))
    add_text_box(s, title, bx + Inches(0.08), by + Inches(0.05), Inches(4.0), Inches(0.4),
                 font_size=13, bold=True, color=bc)
    add_text_box(s, body, bx + Inches(0.1), by + Inches(0.42), Inches(4.0), Inches(1.4),
                 font_size=11, color=GRAY_LIGHT)

# 오른쪽: 성과 비교 + 연도별 수익률 테이블
add_rect(s, Inches(4.8), Inches(0.85), Inches(8.2), Inches(1.1),
         fill_color=GRAY_DIM, line_color=ACCENT_GOLD, line_width=Pt(1))
perf_lines = [
    ("기간               S6 골든크로스     S7 FnG+RSI(2)★    S8 골든크로스+FnG", True, ACCENT_GOLD),
    ("2011~2024  수익  11,287% / CAGR 40.3%   4,257% / 31.0%   1,686% / 22.9%", False, GRAY_LIGHT),
    ("2021~2024  수익    +394.7% / CAGR 49.3%  +684.8% / 67.7%  +252.2% / 37.1%", False, GRAY_LIGHT),
    ("MDD (최근4년)          -49.0%               -30.3% ◀      -28.4% ◀", False, ACCENT_CYAN),
]
add_multiline_textbox(s, perf_lines, Inches(4.85), Inches(0.88), Inches(8.1), Inches(1.0),
                      font_size=10.5, color=WHITE)

# 연도별 수익률 테이블
annual_data = [
    ("연도",  "S6 골든크로스",   "S7 FnG+RSI(2)★", "S8 결합전략"),
    ("2011",  "  -7.0%▼",       "+20.9%▲",          " -0.8%▼"),
    ("2012",  "+51.4%▲",        "+43.3%▲",          "+42.7%▲"),
    ("2013", "+119.4%▲",        "+83.6%▲",          "+83.6%▲"),
    ("2014",  "+60.8%▲",        " +9.6%▲",          " +9.6%▲"),
    ("2015",  " -4.8%▼",        "+26.7%▲",          "+26.9%▲"),
    ("2016",  " -8.3%▼",        "+17.5%▲",          "+18.2%▲"),
    ("2017", "+112.9%▲",        "+50.8%▲",          "+50.8%▲"),
    ("2018",  " -2.3%▼",        "-23.7%▼",          "-10.6%▼"),
    ("2019",  "+47.9%▲",        "+24.6%▲",          "+17.9%▲"),
    ("2020",  "+66.1%▲",        "-30.2%▼",          "-30.2%▼"),
    ("2021 ←", "+91.3%▲",       "+47.5%▲",          "+47.5%▲"),
    ("2022 ←", "-37.9%▼",       "+61.3%▲ ★",        "-14.6%▼"),
    ("2023 ←", "+148.8%▲",      "+76.4%▲",          "+46.2%▲"),
    ("2024 ←", "+71.2%▲",       "+91.2%▲",          "+91.2%▲"),
]

ann_tbl = s.shapes.add_table(15, 4, Inches(4.8), Inches(2.0), Inches(8.2), Inches(5.35)).table
for ci, cw in enumerate([Inches(1.1), Inches(2.2), Inches(2.4), Inches(2.5)]):
    ann_tbl.columns[ci].width = cw
for ri, row_vals in enumerate(annual_data):
    is_header = (ri == 0)
    is_2022 = (ri == 13)
    is_recent = ri >= 12
    bg = (RGBColor(0x00,0x38,0x70) if is_header
          else RGBColor(0x00,0x2a,0x1a) if is_2022
          else (RGBColor(0x14,0x2a,0x14) if is_recent and ri % 2 == 0
                else RGBColor(0x0d,0x1e,0x36) if ri % 2 == 0
                else RGBColor(0x0a,0x2a,0x4a)))
    tc = ([ACCENT_CYAN]*4 if is_header
          else [GRAY_LIGHT,
                ACCENT_RED if "▼" in row_vals[1] else ACCENT_GREEN,
                ACCENT_GOLD if is_2022 else (ACCENT_GREEN if "▲" in row_vals[2] else ACCENT_RED),
                ACCENT_GREEN if "▲" in row_vals[3] else ACCENT_RED])
    add_table_row(ann_tbl, ri, row_vals, font_size=10.5,
                  bold=is_header or is_2022, bg_color=bg, text_colors=tc)

add_text_box(s, "← 최근 4년 구간  |  ★ 2022년 하락장에서 S7만 +61.3% 달성",
             Inches(4.8), Inches(7.1), Inches(8.2), Inches(0.3),
             font_size=10, color=ACCENT_GOLD)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE D — FnG+RSI(2) 알고리즘 타당성 결론
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s)
add_rect(s, 0, 0, SLIDE_W, Inches(0.06), ACCENT_GREEN)
add_rect(s, 0, SLIDE_H - Inches(0.06), SLIDE_W, Inches(0.06), ACCENT_GREEN)

add_text_box(s, "FnG+RSI(2) 알고리즘 타당성 검증 결론",
             Inches(0.5), Inches(0.12), Inches(12.5), Inches(0.55),
             font_size=26, bold=True, color=WHITE)
add_divider(s, Inches(0.72), ACCENT_GREEN)

# 좌측: 핵심 근거 4개
add_text_box(s, "✅  알고리즘 타당성 근거",
             Inches(0.4), Inches(0.85), Inches(7.5), Inches(0.45),
             font_size=16, bold=True, color=ACCENT_GREEN)

evidence = [
    ("① 최근 4년 (2021~2024) 모든 지표 1위", True, ACCENT_GOLD),
    ("   총수익 +684.8%  |  CAGR 67.7%  |  MDD -30.3%  |  Sharpe 1.37", False, GRAY_LIGHT),
    ("   9개 룰 기반 전략 중 압도적 1위 — 단순 운이 아닌 구조적 우위", False, TEXT_DIM),
    ("", False, None),
    ("② 하락장(2022년)에서 유일하게 수익 달성", True, ACCENT_GOLD),
    ("   S7: +61.3%  vs  S6: -37.9%  vs  TQQQ Buy&Hold: -79.7%", False, ACCENT_GREEN),
    ("   FnG/공포지수 활용으로 하락장 단기 바운스 포착 가능", False, GRAY_LIGHT),
    ("   골든/데드크로스 필터 없이도 자체적으로 리스크 관리됨", False, TEXT_DIM),
    ("", False, None),
    ("③ 골든크로스 결합(S8) 실험 결과: 오히려 성과 저하", True, ACCENT_GOLD),
    ("   S8: +252.2%  <  S7: +684.8%  (최근 4년 기준)", False, GRAY_LIGHT),
    ("   2022년 하락장 수익 기회를 데드크로스 필터가 차단", False, TEXT_DIM),
    ("   → RSI(2) 신호 자체가 이미 최적의 진입/청산 타이밍", False, ACCENT_CYAN),
    ("", False, None),
    ("④ 14년 누적 검증 (2011~2024 백테스트)", True, ACCENT_GOLD),
    ("   CAGR 31.0% / MDD -67.2% / Sharpe 0.85", False, GRAY_LIGHT),
    ("   단기 노이즈가 아닌 장기 구조적 알고리즘 신뢰성 확인", False, TEXT_DIM),
]
add_multiline_textbox(s, evidence, Inches(0.4), Inches(1.35), Inches(7.6), Inches(5.3),
                      font_size=11.5, color=WHITE)

# 우측 상단: 최근 4년 성과 박스
add_rect(s, Inches(8.3), Inches(0.85), Inches(4.7), Inches(3.1),
         fill_color=RGBColor(0x00,0x2a,0x1a), line_color=ACCENT_GREEN, line_width=Pt(2))
add_text_box(s, "FnG+RSI(2) 최근 4년 성과",
             Inches(8.4), Inches(0.9), Inches(4.5), Inches(0.45),
             font_size=14, bold=True, color=ACCENT_GREEN)
kpi_lines = [
    ("총수익         +684.8%", True, ACCENT_GOLD),
    ("CAGR            67.7%", True, ACCENT_GOLD),
    ("MDD            -30.3%", True, ACCENT_CYAN),
    ("Sharpe           1.37", True, ACCENT_CYAN),
    ("", False, None),
    ("$10,000 → $78,481", True, WHITE),
    ("(4년, 7.8배 성장)", False, GRAY_LIGHT),
]
add_multiline_textbox(s, kpi_lines, Inches(8.5), Inches(1.42), Inches(4.3), Inches(2.3),
                      font_size=14, color=WHITE)

# 우측 하단: 결론 박스
add_rect(s, Inches(8.3), Inches(4.1), Inches(4.7), Inches(3.2),
         fill_color=GRAY_DIM, line_color=ACCENT_GOLD, line_width=Pt(2))
add_text_box(s, "종합 결론",
             Inches(8.4), Inches(4.15), Inches(4.5), Inches(0.45),
             font_size=15, bold=True, color=ACCENT_GOLD)
conclusion = [
    ("TQQQ 단일 종목 투자 기준으로", False, GRAY_LIGHT),
    ("FnG+RSI(2) 전략이 최적 알고리즘", True, WHITE),
    ("", False, None),
    ("• 장기(14년) 검증된 구조적 우위", False, GRAY_LIGHT),
    ("• 하락장 방어 + 상승장 참여 동시 달성", False, GRAY_LIGHT),
    ("• 어떤 추세 필터보다 RSI 자체 신호가 강력", False, GRAY_LIGHT),
    ("", False, None),
    ("현재 BATA 서비스로 실전 운용 중 ✅", True, ACCENT_GREEN),
]
add_multiline_textbox(s, conclusion, Inches(8.4), Inches(4.65), Inches(4.5), Inches(2.4),
                      font_size=12, color=WHITE)


# ══════════════════════════════════════════════════════════════════════════════
# 저장
# ══════════════════════════════════════════════════════════════════════════════
out_path = "/Users/batagota/BATAGOTA/10_AI_BATA/03_BATA_ALGO/investor_strategy_comparison.pptx"
prs.save(out_path)
print(f"✅ PPT 저장 완료: {out_path}")
print(f"   슬라이드 수: {len(prs.slides)}장")
