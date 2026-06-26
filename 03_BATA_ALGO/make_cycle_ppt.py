"""
make_cycle_ppt.py
──────────────────────────────────────────────────────────────
RSI(2)+FnG 사이클 완료 시 자동 PPT 생성 및 이메일 발송

메모리 분리 설계:
  - create_rsi_fng_sheet.py / sheets_manager / Google Sheets 등
    무거운 모듈은 일절 임포트하지 않음
  - RSI 수학 함수 및 시뮬레이션 로직을 이 파일에 인라인으로 포함
  - 사이클 데이터 추출 후 전체 df 를 즉시 해제

사용법:
  # 마지막 완료 TQQQ 사이클 매도종료 PPT
  python make_cycle_ppt.py

  # 종목 지정
  python make_cycle_ppt.py SOXL

  # 앱 비밀번호 환경변수 사용
  APP_PASS=xxxx python make_cycle_ppt.py
"""

from __future__ import annotations

import os
import smtplib
import sys
from datetime import datetime, timedelta
from email import encoders
from email.mime.base import MIMEBase
from email.mime.multipart import MIMEMultipart
from email.mime.text import MIMEText
from pathlib import Path

# .env 파일 자동 로드 (python-dotenv 없이 직접 파싱)
def _load_env(env_path: Path):
    if not env_path.exists():
        return
    for line in env_path.read_text(encoding="utf-8").splitlines():
        line = line.strip()
        if not line or line.startswith("#") or "=" not in line:
            continue
        key, _, val = line.partition("=")
        key = key.strip()
        val = val.strip()
        if key and key not in os.environ:
            os.environ[key] = val

_load_env(Path(__file__).parent / ".env")

import numpy as np
import pandas as pd
import yfinance as yf
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

# fear_greed_history 는 같은 디렉터리에 있음 (pandas + curl 만 사용, 경량)
_ALGO_DIR = Path(__file__).parent
sys.path.insert(0, str(_ALGO_DIR))
from fear_greed_history import load_fng_history  # noqa: E402

# ─────────────────────────────────────────────────────────────
# 설정 상수
# ─────────────────────────────────────────────────────────────
TICKER_CONFIG: dict[str, dict] = {
    "TQQQ": {"period": 2, "buy_below": 15.0, "sell_above": 75.0},
    "SOXL": {"period": 2, "buy_below": 15.0, "sell_above": 90.0},
}
FEAR_MAX  = 25
GREED_MIN = 75
CAPITAL              = 100_000.0
ANNUAL_START_CAPITAL = 100_000.0   # 보고서 표시용 연도별 시작 투입금 ($)
START_DATE = "2021-01-01"

SENDER   = "hongryong.cho@gmail.com"
RECEIVER = "hongryong.cho@gmail.com"

OUT_DIR = _ALGO_DIR / "cycle_ppt_output"

# ─────────────────────────────────────────────────────────────
# 색상 팔레트 (make_investor_ppt.py 동일)
# ─────────────────────────────────────────────────────────────
BG_DARK      = RGBColor(0x1e, 0x2e, 0x52)   # 슬라이드 배경 (미디엄 네이비)
ACCENT_BLUE  = RGBColor(0x00, 0x9f, 0xff)
ACCENT_CYAN  = RGBColor(0x00, 0xe5, 0xff)
ACCENT_GREEN = RGBColor(0x00, 0xd4, 0x8a)
ACCENT_RED   = RGBColor(0xff, 0x4d, 0x6d)
ACCENT_GOLD  = RGBColor(0xff, 0xd7, 0x00)
WHITE        = RGBColor(0xf0, 0xf4, 0xff)
GRAY_LIGHT   = RGBColor(0xc0, 0xca, 0xe0)
GRAY_MID     = RGBColor(0x38, 0x50, 0x7c)   # 테이블 헤더 배경
GRAY_DIM     = RGBColor(0x2a, 0x3c, 0x64)   # 테이블 데이터 행 / 지표 박스
TEXT_DIM     = RGBColor(0x88, 0x9e, 0xbc)
FONT_NAME    = "Malgun Gothic"

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ─────────────────────────────────────────────────────────────
# RSI(2) + LOC 시뮬레이션 — 인라인 (외부 모듈 미임포트)
# ─────────────────────────────────────────────────────────────

def _solve_rsi_target_close(prev: float, mu: float, md: float,
                             target_rsi: float, alpha: float) -> float:
    """Wilder EWM 상태에서 target_rsi 를 만드는 다음 종가를 역산"""
    if mu == 0.0 and md == 0.0:
        return round(prev, 2)
    beta     = 1.0 - alpha
    r_target = target_rsi / (100.0 - target_rsi)
    x_down   = prev + (md - mu / r_target) * beta / alpha
    if x_down <= prev:
        return round(x_down, 2)
    x_up = prev + (r_target * md - mu) * beta / alpha
    return round(x_up, 2)


def _compute_limit_prices(closes, mu_arr, md_arr, buy_below, sell_above, alpha):
    n = len(closes)
    buy_lims  = np.full(n, np.nan)
    sell_lims = np.full(n, np.nan)
    for i in range(1, n):
        buy_lims[i]  = _solve_rsi_target_close(closes[i-1], mu_arr[i-1], md_arr[i-1], buy_below,  alpha)
        sell_lims[i] = _solve_rsi_target_close(closes[i-1], mu_arr[i-1], md_arr[i-1], sell_above, alpha)
    return buy_lims, sell_lims


def _compute_next_day_limit_prices(closes, mu_arr, md_arr, buy_below, sell_above, alpha):
    n = len(closes)
    nbl = np.full(n, np.nan)
    nsl = np.full(n, np.nan)
    for i in range(n):
        nbl[i] = _solve_rsi_target_close(closes[i], mu_arr[i], md_arr[i], buy_below,  alpha)
        nsl[i] = _solve_rsi_target_close(closes[i], mu_arr[i], md_arr[i], sell_above, alpha)
    return nbl, nsl


def simulate_fng(
    close: pd.Series,
    fng: pd.Series,
    period: int,
    buy_below: float,
    sell_above: float,
    fear_max: int = FEAR_MAX,
    greed_min: int = GREED_MIN,
) -> pd.DataFrame:
    """
    RSI(2)+FnG 시뮬레이션. create_rsi_fng_sheet.simulate_with_fng 과 동일한 로직,
    양도세 추적 제외 (사이클 날짜/가격/RSI/F&G 추출 목적으로만 사용).

    Returns: DataFrame with columns:
        date, close, chg_pct, fng, rsi2, fng_blocked, action,
        holdings, cash, total_assets, buy_limit_px, sell_limit_px
    """
    alpha  = 1.0 / period
    closes = close.values.astype(float)
    n      = len(closes)

    # Wilder EWM
    mu = np.zeros(n)
    md = np.zeros(n)
    for i in range(1, n):
        delta  = closes[i] - closes[i - 1]
        mu[i]  = alpha * max(delta, 0.0)  + (1 - alpha) * mu[i - 1]
        md[i]  = alpha * max(-delta, 0.0) + (1 - alpha) * md[i - 1]

    with np.errstate(divide="ignore", invalid="ignore"):
        rs  = np.where(md == 0, np.inf, mu / md)
        rsi = 100.0 - 100.0 / (1.0 + rs)
    rsi[0] = np.nan

    _, _   = _compute_limit_prices(closes, mu, md, buy_below, sell_above, alpha)
    nbl, nsl = _compute_next_day_limit_prices(closes, mu, md, buy_below, sell_above, alpha)

    fng_al = fng.reindex(close.index, method="ffill").fillna(50).astype(int)

    cash     = CAPITAL
    holdings = 0.0
    rows     = []
    cash_at_buy = 0.0

    for i, (dt, price) in enumerate(close.items()):
        fng_val  = int(fng_al.iloc[i])
        action   = "HOLD"
        blocked  = ""
        prev_nbl = nbl[i - 1] if i > 0 else np.nan
        prev_nsl = nsl[i - 1] if i > 0 else np.nan

        if i > 0:
            pos_cash = holdings == 0.0
            if pos_cash and not np.isnan(prev_nbl) and float(price) <= float(prev_nbl):
                if fng_val >= greed_min:
                    blocked = f"BUY차단(F&G={fng_val}≥{greed_min})"
                else:
                    cash_at_buy = cash
                    holdings    = cash / price
                    cash        = 0.0
                    action      = "BUY"
            elif (not pos_cash) and not np.isnan(prev_nsl) and float(price) >= float(prev_nsl):
                if fng_val <= fear_max:
                    blocked = f"SELL차단(F&G={fng_val}≤{fear_max})"
                else:
                    cash      = holdings * price
                    holdings  = 0.0
                    action    = "SELL"

        total = cash + holdings * float(price)
        prev_price = closes[i - 1] if i > 0 else float(price)
        chg = round((float(price) / prev_price - 1) * 100, 2) if i > 0 else 0.0

        rows.append({
            "date":         dt.strftime("%Y-%m-%d"),
            "close":        round(float(price), 2),
            "chg_pct":      chg,
            "fng":          fng_val,
            "rsi2":         round(float(rsi[i]), 2) if not np.isnan(rsi[i]) else None,
            "fng_blocked":  blocked,
            "action":       action,
            "holdings":     round(float(holdings), 6),
            "cash":         round(float(cash), 2),
            "total_assets": round(float(total), 2),
            "buy_limit_px":  round(float(prev_nbl), 2) if not np.isnan(prev_nbl) else None,
            "sell_limit_px": round(float(prev_nsl), 2) if not np.isnan(prev_nsl) else None,
        })

    df = pd.DataFrame(rows)
    df.index = close.index
    return df


def extract_cycles(df: pd.DataFrame) -> tuple[list[dict], dict | None]:
    """_extract_cycle_state_fng 인라인 버전 (양도세 없음)"""
    completed   = []
    current_buy = None
    cycle_no    = 0

    for row in df.itertuples(index=False):
        if row.action == "BUY":
            cycle_no += 1
            current_buy = {
                "cycle_no":  cycle_no,
                "buy_date":  row.date,
                "buy_price": row.close,
                "buy_rsi":   row.rsi2,
                "buy_fng":   row.fng,
                "start_cash": CAPITAL if cycle_no == 1 else row.cash,  # approximate
            }
        elif row.action == "SELL" and current_buy is not None:
            start = current_buy.get("start_cash") or CAPITAL
            # start_cash: cash before buy (나중에 df에서 직접 읽음)
            completed.append({
                **current_buy,
                "sell_date":  row.date,
                "sell_price": row.close,
                "sell_rsi":   row.rsi2,
                "sell_fng":   row.fng,
                "end_cash":   row.cash,
            })
            current_buy = None

    # start_cash 를 df 에서 직접 추출 (BUY 직전 행의 total_assets)
    rows_list = list(df.itertuples(index=False))
    for cyc in completed:
        buy_idx = next((i for i, r in enumerate(rows_list) if r.date == cyc["buy_date"]), None)
        if buy_idx and buy_idx > 0:
            cyc["start_cash"] = rows_list[buy_idx - 1].total_assets
        else:
            cyc["start_cash"] = CAPITAL
        start = cyc["start_cash"]
        end   = cyc["end_cash"]
        cyc["cycle_return_pct"]    = round((end / start - 1) * 100, 2)
        cyc["cycle_return_amount"] = round(end - start, 2)
        cyc["days_held"] = (
            pd.to_datetime(cyc["sell_date"]) - pd.to_datetime(cyc["buy_date"])
        ).days

    open_cycle = None
    if current_buy is not None:
        last = df.iloc[-1]
        start = current_buy.get("start_cash") or CAPITAL
        open_cycle = {
            **current_buy,
            "sell_date":  "진행중",
            "sell_price": float(last["close"]),
            "sell_rsi":   None,
            "sell_fng":   int(last["fng"]),
            "end_cash":   float(last["total_assets"]),
            "cycle_return_pct": round((float(last["total_assets"]) / start - 1) * 100, 2),
            "cycle_return_amount": round(float(last["total_assets"]) - start, 2),
            "days_held": (
                pd.to_datetime(last["date"]) - pd.to_datetime(current_buy["buy_date"])
            ).days,
        }

    return completed, open_cycle


# ─────────────────────────────────────────────────────────────
# 연도별 사이클 번호 및 복리 투자금 계산
# ─────────────────────────────────────────────────────────────

def _compute_year_info(cycle: dict, completed_cycles: list) -> dict:
    """
    연도별 YY-N 라벨과 복리 투입/정산금을 계산한다.
    백테스트 start_cash / end_cash 를 절대 사용하지 않음.

    반환:
      year_label      : "26-3" 형식
      year_cycle_no   : 연도 내 순번 (1-based)
      buy_year        : 4자리 연도 (int)
      start_amount    : 이 사이클 시작 시점의 자산 ($)
      end_amount      : 이 사이클 종료 시점의 자산 ($)
      profit          : 이 사이클 수익금 ($)
      ytd_pct_before  : 이 사이클 시작 전 YTD 누적 수익률 (매수시작 커버용)
      ytd_pct_after   : 이 사이클 종료 후 YTD 누적 수익률 (매도종료 커버용)
      ytd_end_before  : 이 사이클 시작 시점 자산 = start_amount
      ytd_end_after   : 이 사이클 종료 시점 자산 = end_amount
    """
    buy_year = int(str(cycle["buy_date"])[:4])
    yy = str(buy_year)[-2:]

    # 같은 연도 사이클만 time-sorted로 추출
    year_cycles = [c for c in completed_cycles
                   if int(str(c["buy_date"])[:4]) == buy_year]

    # 현재 사이클이 연도 내 몇 번째?
    year_cycle_no = 1
    for i, c in enumerate(year_cycles):
        if c["cycle_no"] == cycle["cycle_no"]:
            year_cycle_no = i + 1
            break

    year_label = f"{yy}-{year_cycle_no}"

    # 복리 누적: 연도 시작 $100K → 현 사이클까지
    amount = ANNUAL_START_CAPITAL
    start_amount = ANNUAL_START_CAPITAL
    for c in year_cycles:
        start_amount = amount
        amount = amount * (1 + c["cycle_return_pct"] / 100.0)
        if c["cycle_no"] == cycle["cycle_no"]:
            break

    end_amount = amount

    ytd_pct_before = (start_amount - ANNUAL_START_CAPITAL) / ANNUAL_START_CAPITAL * 100
    ytd_pct_after  = (end_amount  - ANNUAL_START_CAPITAL) / ANNUAL_START_CAPITAL * 100

    return {
        "year_label":     year_label,
        "year_cycle_no":  year_cycle_no,
        "buy_year":       buy_year,
        "start_amount":   round(start_amount, 2),
        "end_amount":     round(end_amount, 2),
        "profit":         round(end_amount - start_amount, 2),
        "ytd_pct_before": round(ytd_pct_before, 2),
        "ytd_pct_after":  round(ytd_pct_after, 2),
        "ytd_end_before": round(start_amount, 2),
        "ytd_end_after":  round(end_amount, 2),
    }


def _year_return_stats(completed_cycles: list) -> dict:
    """연도별 복리 수익률·정산금액·MDD (보고서 표시용, 매년 ANNUAL_START_CAPITAL 재시작)."""
    year_data: dict[int, list[float]] = {}
    for c in completed_cycles:
        y = int(str(c["buy_date"])[:4])
        year_data.setdefault(y, []).append(c["cycle_return_pct"])

    result = {}
    for year in sorted(year_data):
        amt  = ANNUAL_START_CAPITAL
        peak = amt
        mdd  = 0.0
        for r in year_data[year]:
            amt = amt * (1 + r / 100.0)
            if amt > peak:
                peak = amt
            dd = (peak - amt) / peak * 100.0
            if dd > mdd:
                mdd = dd
        result[year] = {
            "cycle_count":    len(year_data[year]),
            "annual_ret_pct": round((amt - ANNUAL_START_CAPITAL) / ANNUAL_START_CAPITAL * 100.0, 2),
            "end_amount":     round(amt, 2),
            "mdd_pct":        round(mdd, 2),
        }
    return result


# ─────────────────────────────────────────────────────────────
# PPT 헬퍼 함수
# ─────────────────────────────────────────────────────────────

def _set_bg(slide, color: RGBColor = BG_DARK):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_rect(slide, left, top, width, height,
              fill_color=None, line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def _add_text(slide, text, left, top, width, height,
              size=18, bold=False, color=WHITE,
              align=PP_ALIGN.LEFT, italic=False):
    txb = slide.shapes.add_textbox(left, top, width, height)
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size     = Pt(size)
    run.font.bold     = bold
    run.font.italic   = italic
    run.font.color.rgb = color
    run.font.name     = FONT_NAME
    return txb


def _add_divider(slide, top, color=ACCENT_BLUE, thickness=Pt(1.5)):
    shape = slide.shapes.add_shape(1, Inches(0.4), top,
                                   SLIDE_W - Inches(0.8), thickness)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _add_table(slide, headers, rows_data, left, top, width, height,
               header_bg=GRAY_MID, row_bg=GRAY_DIM, alt_bg=None,
               font_size=11, col_widths=None):
    """범용 테이블 추가. rows_data: list of (values, optional row_bg_color)"""
    n_rows = 1 + len(rows_data)
    n_cols = len(headers)
    tbl    = slide.shapes.add_table(n_rows, n_cols, left, top, width, height).table

    # 컬럼 너비 설정
    if col_widths:
        for ci, w in enumerate(col_widths):
            tbl.columns[ci].width = w

    # 헤더 행
    for ci, h in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_bg
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf = cell.text_frame
        tf.word_wrap = False
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        runs = tf.paragraphs[0].runs
        if runs:
            runs[0].font.size  = Pt(font_size)
            runs[0].font.bold  = True
            runs[0].font.color.rgb = ACCENT_CYAN
            runs[0].font.name  = FONT_NAME

    # 데이터 행
    for ri, row_item in enumerate(rows_data):
        if isinstance(row_item, (list, tuple)) and row_item and isinstance(row_item[-1], RGBColor):
            values, bg = row_item[:-1], row_item[-1]
        elif isinstance(row_item, dict):
            values = row_item.get("values", [])
            bg     = row_item.get("bg", None)
        else:
            values = row_item
            bg = alt_bg if (alt_bg and ri % 2 == 1) else row_bg

        bg = bg or row_bg
        for ci, val in enumerate(values):
            cell = tbl.cell(ri + 1, ci)
            text = str(val) if val is not None else ""
            cell.text = text
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf = cell.text_frame
            tf.word_wrap = False
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            runs = tf.paragraphs[0].runs
            if runs:
                runs[0].font.size  = Pt(font_size)
                runs[0].font.bold  = False
                runs[0].font.color.rgb = WHITE
                runs[0].font.name  = FONT_NAME

    return tbl


def _color_for_return(pct: float) -> RGBColor:
    if pct > 0:
        return ACCENT_GREEN
    elif pct < 0:
        return ACCENT_RED
    return GRAY_LIGHT


def _color_for_action(action: str) -> RGBColor:
    if action == "BUY":
        return ACCENT_GREEN
    if action == "SELL":
        return ACCENT_RED
    if "차단" in action:
        return ACCENT_GOLD
    return GRAY_LIGHT


# ─────────────────────────────────────────────────────────────
# PPT 슬라이드 생성
# ─────────────────────────────────────────────────────────────

def _slide_cover(prs, ticker: str, cycle: dict, ppt_type: str, year_info: dict):
    """슬라이드 1: 커버"""
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    _set_bg(s)

    # 상단/하단 액센트 바
    _add_rect(s, 0, 0, SLIDE_W, Inches(0.10), ACCENT_BLUE)
    _add_rect(s, 0, SLIDE_H - Inches(0.10), SLIDE_W, Inches(0.10), ACCENT_BLUE)

    # 왼쪽 세로 바
    _add_rect(s, 0, Inches(0.10), Inches(0.08), SLIDE_H - Inches(0.20), ACCENT_BLUE)

    # 메인 박스 (YTD 구역 포함하여 조금 더 높게)
    _add_rect(s, Inches(0.5), Inches(1.6), Inches(12.4), Inches(4.8),
              fill_color=GRAY_DIM, line_color=ACCENT_BLUE, line_width=Pt(1.5))

    # 종목 태그
    _add_text(s, ticker, Inches(0.6), Inches(1.7), Inches(2), Inches(0.7),
              size=36, bold=True, color=ACCENT_GOLD, align=PP_ALIGN.CENTER)

    # 타입 뱃지
    type_color = ACCENT_GREEN if "매수" in ppt_type else (
        ACCENT_RED if "매도" in ppt_type else ACCENT_GOLD
    )
    _add_rect(s, Inches(2.8), Inches(1.75), Inches(3.0), Inches(0.55),
              fill_color=type_color)
    _add_text(s, ppt_type, Inches(2.8), Inches(1.75), Inches(3.0), Inches(0.55),
              size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # 사이클 번호 (글로벌 번호 + 연도별 번호)
    _add_text(s, f"Cycle #{cycle['cycle_no']:02d}  ·  {year_info['year_label']}",
              Inches(0.6), Inches(2.5), Inches(7), Inches(0.7),
              size=28, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    # 메인 타이틀
    _add_text(s, f"RSI(2) + Fear & Greed 전략",
              Inches(0.6), Inches(3.1), Inches(12), Inches(0.9),
              size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # 서브타이틀: 날짜 범위
    buy_lbl  = cycle["buy_date"]
    sell_lbl = cycle["sell_date"] if cycle["sell_date"] != "진행중" else "진행중"
    _add_text(s, f"매수: {buy_lbl}  →  매도: {sell_lbl}",
              Inches(0.6), Inches(4.0), Inches(12), Inches(0.6),
              size=20, bold=False, color=GRAY_LIGHT, align=PP_ALIGN.CENTER)

    # 이 사이클 수익률 뱃지
    ret_pct    = cycle.get("cycle_return_pct", 0.0)
    ret_color  = _color_for_return(ret_pct)
    ret_sign   = "+" if ret_pct >= 0 else ""
    _add_rect(s, Inches(4.8), Inches(4.7), Inches(3.7), Inches(0.8), fill_color=ret_color)
    _add_text(s, f"이 사이클  {ret_sign}{ret_pct:.2f}%",
              Inches(4.8), Inches(4.7), Inches(3.7), Inches(0.8),
              size=26, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # 보유기간
    days = cycle.get("days_held", "")
    _add_text(s, f"보유기간  {days}일",
              Inches(0.6), Inches(4.85), Inches(4.0), Inches(0.4),
              size=16, bold=False, color=TEXT_DIM, align=PP_ALIGN.LEFT)

    # ── 연도 누적 수익률 구역 ──────────────────────────────────────
    _add_rect(s, Inches(0.6), Inches(5.55), Inches(12.2), Inches(0.015),
              fill_color=GRAY_MID)
    # 매수시작: 이전 사이클까지의 누적 / 매도종료: 이 사이클 포함 누적
    if "매수" in ppt_type:
        ytd_pct = year_info["ytd_pct_before"]
        ytd_end = year_info["ytd_end_before"]
        ytd_label = f"{year_info['buy_year']}년 기존 누적 (이전 사이클)"
    else:
        ytd_pct = year_info["ytd_pct_after"]
        ytd_end = year_info["ytd_end_after"]
        ytd_label = f"{year_info['buy_year']}년 누적 수익률"

    ytd_sign  = "+" if ytd_pct >= 0 else ""
    ytd_color = _color_for_return(ytd_pct)
    _add_text(s,
              f"{ytd_label}  {ytd_sign}{ytd_pct:.2f}%"
              f"  ($100,000 → ${ytd_end:,.0f})",
              Inches(0.6), Inches(5.65), Inches(12.2), Inches(0.5),
              size=17, bold=True, color=ytd_color, align=PP_ALIGN.CENTER)

    # 소자막 채널명
    _add_text(s, "BATA  ·  RSI(2) Algorithm Channel",
              Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.35),
              size=13, bold=False, color=TEXT_DIM, align=PP_ALIGN.RIGHT)


def _slide_summary(prs, ticker: str, cycle: dict, cfg: dict,
                   year_info: dict, ppt_type: str):
    """슬라이드 2: 사이클 요약 (주요 지표)"""
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    _set_bg(s)
    _add_rect(s, 0, 0, SLIDE_W, Inches(0.08), ACCENT_BLUE)

    # 타이틀 (연도별 번호 포함)
    _add_text(s,
              f"[{ticker}  {year_info['year_label']}  /  Cycle #{cycle['cycle_no']:02d}]  사이클 요약",
              Inches(0.5), Inches(0.12), Inches(12), Inches(0.65),
              size=24, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    _add_divider(s, Inches(0.82))

    ret_pct   = cycle.get("cycle_return_pct", 0.0)
    ret_color = _color_for_return(ret_pct)
    ret_sign  = "+" if ret_pct >= 0 else ""

    # ── 4대 핵심 지표 박스 ──────────────────────────────────────
    metrics = [
        ("수익률",    f"{ret_sign}{ret_pct:.2f}%",             ret_color),
        ("보유기간",  f"{cycle.get('days_held', '?')}일",       ACCENT_CYAN),
        ("매수가",    f"${cycle['buy_price']:.2f}",             ACCENT_BLUE),
        ("매도가",    f"${cycle.get('sell_price', '-')}",       ACCENT_BLUE),
    ]
    box_w = Inches(2.9)
    box_h = Inches(1.7)
    for idx, (label, value, color) in enumerate(metrics):
        lx = Inches(0.4) + idx * (box_w + Inches(0.15))
        _add_rect(s, lx, Inches(1.0), box_w, box_h, fill_color=GRAY_DIM,
                  line_color=color, line_width=Pt(1.5))
        _add_text(s, label, lx, Inches(1.05), box_w, Inches(0.5),
                  size=14, bold=False, color=TEXT_DIM, align=PP_ALIGN.CENTER)
        _add_text(s, str(value), lx, Inches(1.55), box_w, Inches(0.9),
                  size=28, bold=True, color=color, align=PP_ALIGN.CENTER)

    # ── 상세 지표 (투입금/정산금은 보고서 전용 복리 계산치) ──────
    _add_text(s, "상세 지표", Inches(0.5), Inches(2.9), Inches(5), Inches(0.45),
              size=16, bold=True, color=ACCENT_CYAN)
    _add_divider(s, Inches(3.35), color=GRAY_MID, thickness=Pt(1))

    # 매수시작일 때는 아직 종료 미확정 → 정산금/수익금 미표시
    is_buy_start = "매수" in ppt_type
    end_str    = "진행중" if is_buy_start else f"${year_info['end_amount']:,.0f}"
    profit_val = year_info["profit"]
    profit_str = "-" if is_buy_start else f"{ret_sign}${abs(profit_val):,.0f}"
    profit_clr = TEXT_DIM if is_buy_start else ret_color

    sell_rsi_str = "-" if is_buy_start else str(cycle.get("sell_rsi") or "-")
    sell_fng_str = "-" if is_buy_start else str(cycle.get("sell_fng") or "-")

    details = [
        ("투입금액",    f"${year_info['start_amount']:,.0f}",  WHITE),
        ("정산금액",    end_str,                               WHITE),
        ("수익금",      profit_str,                            profit_clr),
        ("매수 RSI(2)", str(cycle.get("buy_rsi") or "-"),     ACCENT_GREEN),
        ("매수 F&G",    str(cycle.get("buy_fng") or "-"),     ACCENT_GREEN),
        ("매도 RSI(2)", sell_rsi_str,                         ACCENT_RED),
        ("매도 F&G",    sell_fng_str,                         ACCENT_RED),
        ("매수 조건",    f"RSI(2) ≤ {cfg['buy_below']}",      TEXT_DIM),
        ("매도 조건",    f"RSI(2) ≥ {cfg['sell_above']}",     TEXT_DIM),
        ("F&G 필터",    f"매수보류 ≥{GREED_MIN} / 매도보류 ≤{FEAR_MAX}", TEXT_DIM),
    ]

    col1_x = Inches(0.5)
    col3_x = Inches(7.1)
    row_h  = Inches(0.38)
    top0   = Inches(3.5)

    for idx, (label, value, color) in enumerate(details):
        col_idx = idx % 2
        row_idx = idx // 2
        lx = col1_x if col_idx == 0 else col3_x
        ty = top0 + row_idx * row_h
        _add_text(s, label + " :", lx, ty, Inches(2.8), row_h,
                  size=14, bold=False, color=TEXT_DIM)
        _add_text(s, value, lx + Inches(2.0), ty, Inches(3.2), row_h,
                  size=14, bold=True, color=color)

    # 소수점 거래 / 수수료 안내 (보고서 전용 금액 기준 명기)
    _add_text(s,
              "※ 소수점(분수) 거래 기준 — 잔고 전액 투입 가능  |  증권사 수수료 및 세금 미포함",
              Inches(0.5), Inches(5.65), Inches(12.3), Inches(0.4),
              size=12, bold=False, color=ACCENT_GOLD, align=PP_ALIGN.CENTER)

    # 하단 면책
    _add_text(s, "⚠ 본 자료는 교육 목적이며 투자 권유가 아닙니다.",
              Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.35),
              size=11, bold=False, color=TEXT_DIM, align=PP_ALIGN.CENTER)


def _slide_entry(prs, ticker: str, cycle: dict, df_context: pd.DataFrame):
    """슬라이드 3: 매수 진입 분석 (매수일 포함 전후 7일 RSI 흐름)"""
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    _set_bg(s)
    _add_rect(s, 0, 0, SLIDE_W, Inches(0.08), ACCENT_GREEN)

    _add_text(s, f"[{ticker} Cycle #{cycle['cycle_no']:02d}]  매수 진입 분석",
              Inches(0.5), Inches(0.12), Inches(12), Inches(0.65),
              size=26, bold=True, color=WHITE)
    _add_divider(s, Inches(0.82), color=ACCENT_GREEN)

    # 매수 기준가 설명
    _add_text(s,
              f"매수 신호 조건:  RSI(2) ≤ {cfg_for(ticker)['buy_below']}  (전일 종가 역산 LOC 기준가 이하 체결)",
              Inches(0.5), Inches(0.95), Inches(12), Inches(0.45),
              size=14, bold=False, color=TEXT_DIM)

    buy_date = cycle["buy_date"]
    # 매수일 포함 이전 7 거래일 (맥락)
    idx_list = list(df_context.index)
    buy_pos  = next((i for i, r in enumerate(df_context.itertuples()) if r.date == buy_date), None)
    if buy_pos is None:
        _add_text(s, "매수 데이터 없음", Inches(0.5), Inches(1.5), Inches(12), Inches(1),
                  size=18, color=ACCENT_RED)
        return

    start_pos = max(0, buy_pos - 6)
    sub_df    = df_context.iloc[start_pos: buy_pos + 2]  # +2: 매수 다음날까지

    headers = ["날짜", "종가($)", "등락률(%)", "RSI(2)", "F&G", "매매 기준가($)", "액션"]
    col_ws  = [Inches(1.7), Inches(1.3), Inches(1.2), Inches(1.1), Inches(0.9), Inches(2.2), Inches(2.4)]
    total_w = sum(col_ws)

    rows_data = []
    for row in sub_df.itertuples(index=False):
        is_buy = row.action == "BUY"
        bg = GRAY_MID if is_buy else GRAY_DIM
        chg_str = f"{row.chg_pct:+.2f}%" if row.chg_pct else ""
        rsi_str = f"{row.rsi2:.2f}" if row.rsi2 is not None else "-"
        limit_str = f"${row.buy_limit_px}" if row.buy_limit_px else "-"
        act_str  = "★ BUY" if is_buy else (row.fng_blocked or row.action)
        rows_data.append({
            "values": [row.date, f"${row.close:.2f}", chg_str, rsi_str,
                       str(row.fng), limit_str, act_str],
            "bg": GRAY_MID if is_buy else GRAY_DIM,
        })

    tbl = _add_table(s, headers, rows_data,
                     left=Inches(0.5), top=Inches(1.5),
                     width=total_w, height=Inches(4.5),
                     font_size=12, col_widths=col_ws)

    # BUY 행 강조 처리 (텍스트 색)
    for ri, row in enumerate(sub_df.itertuples(index=False)):
        if row.action == "BUY":
            for ci in range(len(headers)):
                cell = tbl.cell(ri + 1, ci)
                runs = cell.text_frame.paragraphs[0].runs
                if runs:
                    runs[0].font.color.rgb = ACCENT_GREEN
                    runs[0].font.bold = True

    _add_text(s,
              f"✅ {buy_date}  매수 체결  |  종가 ${cycle['buy_price']:.2f}  |  RSI(2) {cycle.get('buy_rsi') or '-'}  |  F&G {cycle.get('buy_fng') or '-'}",
              Inches(0.5), Inches(6.15), Inches(12.3), Inches(0.45),
              size=14, bold=True, color=ACCENT_GREEN, align=PP_ALIGN.CENTER)

    _add_text(s, "⚠ 본 자료는 교육 목적이며 투자 권유가 아닙니다.",
              Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.35),
              size=11, bold=False, color=TEXT_DIM, align=PP_ALIGN.CENTER)


def _slide_holding(prs, ticker: str, cycle: dict, df_cycle: pd.DataFrame):
    """슬라이드 4: 보유 기간 현황 (매수일 ~ 매도일)"""
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    _set_bg(s)
    _add_rect(s, 0, 0, SLIDE_W, Inches(0.08), ACCENT_BLUE)

    _add_text(s, f"[{ticker} Cycle #{cycle['cycle_no']:02d}]  보유 기간 현황",
              Inches(0.5), Inches(0.12), Inches(12), Inches(0.65),
              size=26, bold=True, color=WHITE)
    _add_divider(s, Inches(0.82))

    # 총 거래일 / F&G 차단 횟수
    n_days = len(df_cycle)
    n_blocked = (df_cycle["fng_blocked"] != "").sum()
    _add_text(s,
              f"보유 {n_days}거래일  |  F&G 차단 발생 {n_blocked}회",
              Inches(0.5), Inches(0.95), Inches(8), Inches(0.4),
              size=14, color=TEXT_DIM)

    # 테이블이 너무 길면 앞뒤 일부만 표시
    MAX_ROWS = 18
    if n_days <= MAX_ROWS:
        display_df = df_cycle
        truncated  = False
    else:
        half = MAX_ROWS // 2
        display_df = pd.concat([df_cycle.iloc[:half], df_cycle.iloc[-half:]])
        truncated  = True

    headers = ["날짜", "종가($)", "등락률(%)", "RSI(2)", "F&G", "차단 / 액션"]
    col_ws  = [Inches(1.8), Inches(1.4), Inches(1.3), Inches(1.2), Inches(1.0), Inches(5.2)]
    total_w = sum(col_ws)

    rows_data = []
    prev_slice_end = None
    for ri, row in enumerate(display_df.itertuples(index=False)):
        is_buy  = row.action == "BUY"
        is_sell = row.action == "SELL"
        bg = GRAY_MID if (is_buy or is_sell) else GRAY_DIM
        chg_str  = f"{row.chg_pct:+.2f}%" if row.chg_pct else ""
        rsi_str  = f"{row.rsi2:.2f}" if row.rsi2 is not None else "-"
        act_str  = ("★ BUY" if is_buy else "▼ SELL" if is_sell
                    else (row.fng_blocked or "HOLD"))
        rows_data.append({
            "values": [row.date, f"${row.close:.2f}", chg_str, rsi_str,
                       str(row.fng), act_str],
            "bg": bg,
        })

    if truncated:
        # 중간에 "..." 구분자 행 삽입 (halfway point)
        half = MAX_ROWS // 2
        rows_data.insert(half, {
            "values": ["...", "...", "...", "...", "...", f"(중략 {n_days - MAX_ROWS}일)"],
            "bg": GRAY_DIM,
        })

    tbl = _add_table(s, headers, rows_data,
                     left=Inches(0.5), top=Inches(1.45),
                     width=total_w, height=Inches(5.3),
                     font_size=11, col_widths=col_ws)

    # BUY/SELL 행 강조
    display_rows = list(display_df.itertuples(index=False))
    offset = 0
    for ri, row in enumerate(display_rows):
        actual_ri = ri + offset
        if row.action in ("BUY", "SELL"):
            color = ACCENT_GREEN if row.action == "BUY" else ACCENT_RED
            for ci in range(len(headers)):
                cell = tbl.cell(actual_ri + 1, ci)
                runs = cell.text_frame.paragraphs[0].runs
                if runs:
                    runs[0].font.color.rgb = color
                    runs[0].font.bold = True
        if truncated and ri == MAX_ROWS // 2 - 1:
            offset = 1  # "..." 행 때문에 오프셋

    _add_text(s, "⚠ 본 자료는 교육 목적이며 투자 권유가 아닙니다.",
              Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.35),
              size=11, bold=False, color=TEXT_DIM, align=PP_ALIGN.CENTER)


def _slide_exit(prs, ticker: str, cycle: dict, df_context: pd.DataFrame):
    """슬라이드 5: 매도 신호 분석"""
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    _set_bg(s)
    _add_rect(s, 0, 0, SLIDE_W, Inches(0.08), ACCENT_RED)

    _add_text(s, f"[{ticker} Cycle #{cycle['cycle_no']:02d}]  매도 신호 분석",
              Inches(0.5), Inches(0.12), Inches(12), Inches(0.65),
              size=26, bold=True, color=WHITE)
    _add_divider(s, Inches(0.82), color=ACCENT_RED)

    _add_text(s,
              f"매도 신호 조건:  RSI(2) ≥ {cfg_for(ticker)['sell_above']}  (전일 종가 역산 LOC 기준가 이상 체결)",
              Inches(0.5), Inches(0.95), Inches(12), Inches(0.45),
              size=14, bold=False, color=TEXT_DIM)

    sell_date = cycle.get("sell_date", "")
    if sell_date == "진행중":
        _add_text(s, "현재 사이클 진행 중 — 매도 미완료",
                  Inches(0.5), Inches(1.5), Inches(12), Inches(1),
                  size=20, color=ACCENT_GOLD)
        return

    sell_pos = next((i for i, r in enumerate(df_context.itertuples()) if r.date == sell_date), None)
    if sell_pos is None:
        _add_text(s, "매도 데이터 없음", Inches(0.5), Inches(1.5), Inches(12), Inches(1),
                  size=18, color=ACCENT_RED)
        return

    start_pos = max(0, sell_pos - 6)
    sub_df    = df_context.iloc[start_pos: sell_pos + 2]

    headers = ["날짜", "종가($)", "등락률(%)", "RSI(2)", "F&G", "매매 기준가($)", "액션"]
    col_ws  = [Inches(1.7), Inches(1.3), Inches(1.2), Inches(1.1), Inches(0.9), Inches(2.2), Inches(2.4)]
    total_w = sum(col_ws)

    rows_data = []
    for row in sub_df.itertuples(index=False):
        is_sell = row.action == "SELL"
        chg_str  = f"{row.chg_pct:+.2f}%" if row.chg_pct else ""
        rsi_str  = f"{row.rsi2:.2f}" if row.rsi2 is not None else "-"
        limit_str = f"${row.sell_limit_px}" if row.sell_limit_px else "-"
        act_str  = "▼ SELL" if is_sell else (row.fng_blocked or row.action)
        rows_data.append({
            "values": [row.date, f"${row.close:.2f}", chg_str, rsi_str,
                       str(row.fng), limit_str, act_str],
            "bg": GRAY_MID if is_sell else GRAY_DIM,
        })

    tbl = _add_table(s, headers, rows_data,
                     left=Inches(0.5), top=Inches(1.5),
                     width=total_w, height=Inches(4.5),
                     font_size=12, col_widths=col_ws)

    for ri, row in enumerate(sub_df.itertuples(index=False)):
        if row.action == "SELL":
            for ci in range(len(headers)):
                cell = tbl.cell(ri + 1, ci)
                runs = cell.text_frame.paragraphs[0].runs
                if runs:
                    runs[0].font.color.rgb = ACCENT_RED
                    runs[0].font.bold = True

    ret_sign = "+" if cycle["cycle_return_pct"] >= 0 else ""
    _add_text(s,
              f"✅ {sell_date}  매도 체결  |  종가 ${cycle.get('sell_price', '-')}  |  RSI(2) {cycle.get('sell_rsi') or '-'}  |  F&G {cycle.get('sell_fng') or '-'}  |  수익률 {ret_sign}{cycle['cycle_return_pct']:.2f}%",
              Inches(0.5), Inches(6.15), Inches(12.3), Inches(0.45),
              size=14, bold=True, color=ACCENT_RED, align=PP_ALIGN.CENTER)

    _add_text(s, "⚠ 본 자료는 교육 목적이며 투자 권유가 아닙니다.",
              Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.35),
              size=11, bold=False, color=TEXT_DIM, align=PP_ALIGN.CENTER)


def _slide_recent_cycles(prs, ticker: str, completed_cycles: list):
    """슬라이드: 최근 사이클 상세 요약 테이블 (올해 전체 or 최근 5개)."""
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    _set_bg(s)
    _add_rect(s, 0, 0, SLIDE_W, Inches(0.08), ACCENT_GREEN)

    current_year = datetime.now().year
    year_cycles = [c for c in completed_cycles
                   if int(str(c["buy_date"])[:4]) == current_year]

    if len(year_cycles) >= 2:
        show_cycles = year_cycles
        label = f"{current_year}년 사이클 전체 ({len(show_cycles)}건)"
    else:
        n = min(5, len(completed_cycles))
        show_cycles = completed_cycles[-n:]
        label = f"최근 {len(show_cycles)}개 사이클"

    _add_text(s, f"{ticker}  사이클 요약  —  {label}",
              Inches(0.4), Inches(0.12), Inches(12.5), Inches(0.65),
              size=24, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    _add_divider(s, Inches(0.82), color=ACCENT_GREEN)

    # 연도별 $100,000 복리 누적으로 정산금액 계산
    temp_running: dict = {}
    cycle_amounts: dict = {}
    for c in sorted(completed_cycles, key=lambda x: x["buy_date"]):
        y = int(str(c["buy_date"])[:4])
        if y not in temp_running:
            temp_running[y] = ANNUAL_START_CAPITAL
        start_amt = temp_running[y]
        end_amt = start_amt * (1 + c["cycle_return_pct"] / 100.0)
        cycle_amounts[c["cycle_no"]] = (start_amt, end_amt)
        temp_running[y] = end_amt

    headers = ["사이클", "매수일", "매수가($)", "RSI/F&G (매수)", "매도일", "매도가($)", "보유(일)", "수익률", "정산금액($)"]
    COL_W = [
        Inches(0.80), Inches(1.25), Inches(1.10), Inches(1.55),
        Inches(1.25), Inches(1.10), Inches(0.85), Inches(1.05), Inches(1.55),
    ]

    rows_data = []
    for c in show_cycles:
        cn = c["cycle_no"]
        _, end_a = cycle_amounts.get(cn, (ANNUAL_START_CAPITAL, ANNUAL_START_CAPITAL))
        buy_rsi = str(c.get("buy_rsi") or "-")
        buy_fng = str(c.get("buy_fng") or "-")
        sell_date = c.get("sell_date") or "진행중"
        sell_price_str = f"${c['sell_price']:.2f}" if c.get("sell_price") else "-"
        try:
            b = datetime.strptime(str(c["buy_date"]), "%Y-%m-%d")
            if sell_date not in ("진행중", None, ""):
                d = datetime.strptime(str(sell_date), "%Y-%m-%d")
                hold_str = str((d - b).days + 1)
            else:
                hold_str = "-"
        except Exception:
            hold_str = "-"
        ret = c["cycle_return_pct"]
        ret_str = f"{ret:+.2f}%"
        rows_data.append({
            "values": [
                f"#{cn}", str(c["buy_date"]), f"${c['buy_price']:.2f}",
                f"RSI {buy_rsi}  /  F&G {buy_fng}",
                str(sell_date), sell_price_str,
                hold_str, ret_str, f"${end_a:,.0f}",
            ],
            "bg": GRAY_DIM,
        })

    n_show = len(show_cycles)
    row_h  = Inches(0.50) if n_show <= 7 else Inches(0.42)
    tbl_h  = Inches(0.44) + n_show * row_h
    tbl_w  = sum(COL_W)

    tbl = _add_table(
        s, headers, rows_data,
        left=Inches(0.30), top=Inches(0.98),
        width=tbl_w, height=tbl_h,
        header_bg=GRAY_MID, row_bg=GRAY_DIM,
        font_size=10, col_widths=COL_W,
    )

    # 수익률 컬럼(7번째, 0-indexed) 색상 적용
    for ri, c in enumerate(show_cycles):
        ret = c["cycle_return_pct"]
        clr = ACCENT_GREEN if ret >= 0 else ACCENT_RED
        cell = tbl.cell(ri + 1, 7)
        runs = cell.text_frame.paragraphs[0].runs
        if runs:
            runs[0].font.color.rgb = clr
            runs[0].font.bold = True

    _add_text(
        s,
        "※ 정산금액: 해당 연도 $100,000 기준 복리 누적 · 소수점 거래 · 수수료/세금 미포함",
        Inches(0.35), Inches(7.10), Inches(12.6), Inches(0.32),
        size=10, bold=False, color=TEXT_DIM, align=PP_ALIGN.CENTER,
    )


def _slide_reference(prs, ticker: str, completed_cycles: list):
    """슬라이드: RSI(2)/F&G/MDD 지표 설명 + 연도별 수익률 테이블."""
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    _set_bg(s)
    _add_rect(s, 0, 0, SLIDE_W, Inches(0.08), ACCENT_CYAN)

    _add_text(s, "지표 설명  &  연도별 수익률",
              Inches(0.5), Inches(0.12), Inches(12), Inches(0.65),
              size=26, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    _add_divider(s, Inches(0.82), color=ACCENT_CYAN)

    cfg = cfg_for(ticker)

    # ── 좌측: 지표 설명 박스 3개 ─────────────────────────────────
    LX = Inches(0.4)
    LW = Inches(5.3)

    # 각 라인: (텍스트, 폰트크기, bold, 색상)
    indicators = [
        (
            f"RSI(2)  —  Relative Strength Index",
            [
                ("2일 종가 변화를 Wilder 지수평활(α=0.5)해 0~100으로 환산.", 10, False, GRAY_LIGHT),
                (f"매수 신호 ≤ {cfg['buy_below']}  (과매도)   /   매도 신호 ≥ {cfg['sell_above']}  (과매수)", 10, True, GRAY_LIGHT),
                ("평균상승 = 0.5 × max(오늘-어제, 0) + 0.5 × 전일 평균상승", 9, False, TEXT_DIM),
                ("평균하락 = 0.5 × max(어제-오늘, 0) + 0.5 × 전일 평균하락", 9, False, TEXT_DIM),
                ("RSI(2) = 100 − 100 ÷ (1 + 평균상승 ÷ 평균하락)", 9, True, ACCENT_BLUE),
            ],
            ACCENT_BLUE,
        ),
        (
            "F&G  —  CNN Fear & Greed Index  (0~100)",
            [
                ("CNN이 7가지 시장 데이터를 매일 집계해 0(공포)~100(탐욕)으로 수치화.", 10, False, GRAY_LIGHT),
                (f"매수 보류 ≥ {GREED_MIN} (탐욕)   /   매도 보류 ≤ {FEAR_MAX} (공포)", 10, True, GRAY_LIGHT),
                ("① 모멘텀(S&P500/125일MA)  ② 주가강도(52주 고저)  ③ 거래량추세", 9, False, TEXT_DIM),
                ("④ 풋/콜 비율  ⑤ VIX 변동성  ⑥ 안전자산 수요  ⑦ 정크본드 수요", 9, False, TEXT_DIM),
                ("→ 7가지 지표를 동일 가중 평균해 0~100 단일 지수로 환산", 9, True, ACCENT_GOLD),
            ],
            ACCENT_GOLD,
        ),
        (
            "MDD  —  Maximum Drawdown  (최대 낙폭)",
            [
                ("투자 기간 내 최고 자산 대비 최대 손실률 (낮을수록 전략 안정).", 10, False, GRAY_LIGHT),
                ("연도별: 매년 $100,000 재시작 → 사이클 정산마다 자산값 추적.", 10, False, GRAY_LIGHT),
                ("MDD(%) = (최고정산자산 − 현재정산자산) ÷ 최고정산자산 × 100", 9, False, TEXT_DIM),
                ("최고점: 연도 시작~현재까지 매 사이클 정산 후 누적 최댓값", 9, False, TEXT_DIM),
                ("→ 일중 최저값 아닌 사이클 단위 정산 기준으로 낙폭 측정", 9, True, ACCENT_RED),
            ],
            ACCENT_RED,
        ),
    ]

    ty    = Inches(0.90)
    BOX_H = Inches(2.08)
    GAP   = Inches(0.12)

    for title, lines, color in indicators:
        _add_rect(s, LX, ty, LW, BOX_H, fill_color=GRAY_DIM,
                  line_color=color, line_width=Pt(1.2))
        _add_text(s, title, LX + Inches(0.12), ty + Inches(0.10),
                  LW - Inches(0.24), Inches(0.40),
                  size=13, bold=True, color=color)
        line_y = ty + Inches(0.54)
        for text, sz, bold, clr in lines:
            _add_text(s, text, LX + Inches(0.15), line_y,
                      LW - Inches(0.30), Inches(0.30),
                      size=sz, bold=bold, color=clr)
            line_y += Inches(0.29)
        ty += BOX_H + GAP

    # 세로 구분선
    _add_rect(s, Inches(5.85), Inches(0.9), Inches(0.025), Inches(6.0),
              fill_color=GRAY_MID)

    # ── 우측: 연도별 수익률 테이블 ────────────────────────────────
    RX = Inches(6.1)
    RW = Inches(6.8)

    _add_text(s,
              "연도별 수익률  (연 $100,000 기준 · 소수점 거래 · 수수료/세금 미포함)",
              RX, Inches(0.9), RW, Inches(0.38),
              size=11, bold=False, color=TEXT_DIM, align=PP_ALIGN.CENTER)

    headers = ["연도", "사이클", "연간 수익률", "정산금액", "MDD"]
    COL_W   = [Inches(1.05), Inches(0.9), Inches(1.7), Inches(1.95), Inches(1.1)]

    year_stats   = _year_return_stats(completed_cycles)
    current_year = datetime.now().year
    past_years   = {y: v for y, v in year_stats.items() if y < current_year}

    rows_data = []
    for year, st in sorted(past_years.items()):
        ret_pct  = st["annual_ret_pct"]
        ret_sign = "+" if ret_pct >= 0 else ""
        mdd_str  = f"-{st['mdd_pct']:.1f}%"
        rows_data.append({
            "values": [
                str(year),
                f"{st['cycle_count']}회",
                f"{ret_sign}{ret_pct:.2f}%",
                f"${st['end_amount']:,.0f}",
                mdd_str,
            ],
            "bg": GRAY_DIM,
        })

    n_rows = 1 + len(rows_data)
    tbl_h  = Inches(0.4) + len(rows_data) * Inches(0.48)
    tbl = _add_table(s, headers, rows_data,
                     left=RX, top=Inches(1.35),
                     width=RW, height=tbl_h,
                     header_bg=GRAY_MID, row_bg=GRAY_DIM,
                     font_size=12, col_widths=COL_W)

    # 수익률 칼럼(2) 색상 강조
    for ri, (year, st) in enumerate(sorted(past_years.items())):
        ret_pct = st["annual_ret_pct"]
        color   = _color_for_return(ret_pct)
        cell = tbl.cell(ri + 1, 2)
        runs = cell.text_frame.paragraphs[0].runs
        if runs:
            runs[0].font.color.rgb = color
            runs[0].font.bold = True

    _add_text(s, "⚠ 본 자료는 교육 목적이며 투자 권유가 아닙니다.",
              Inches(0.5), Inches(7.07), Inches(12.3), Inches(0.35),
              size=11, bold=False, color=TEXT_DIM, align=PP_ALIGN.CENTER)


def _make_12m_chart(df: pd.DataFrame, completed_cycles: list,
                    ticker: str) -> bytes:
    """최근 12개월 백테스트 차트 PNG 반환 (단일 패널: 주가 + 자산가치 + 보유기간 음영)."""
    import io
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    import matplotlib.dates as mdates
    import matplotlib.patches as mpatches
    plt.rcParams["font.family"] = "Apple SD Gothic Neo"

    # 최근 12개월 필터
    cutoff = (datetime.now() - timedelta(days=365)).strftime("%Y-%m-%d")
    df12 = df[df["date"] >= cutoff].copy()
    if df12.empty:
        df12 = df.copy()

    # 전체 기간 포트폴리오 계산 후 12개월 구간 값만 추출
    holding_info: dict = {}   # date → (buy_price, start_amount)
    sell_settled: dict = {}   # sell_date → settled value after cycle
    running = ANNUAL_START_CAPITAL
    for c in completed_cycles:
        sa  = running
        bpx = c["buy_price"]
        mask = (df["date"] >= c["buy_date"]) & (df["date"] <= c["sell_date"])
        for d in df.loc[mask, "date"].values:
            holding_info[d] = (bpx, sa)
        running = running * (1 + c["cycle_return_pct"] / 100.0)
        sell_settled[c["sell_date"]] = running

    last_settled = ANNUAL_START_CAPITAL
    port_map: dict = {}
    for _, row in df.iterrows():
        d = row["date"]
        if d in holding_info:
            bp, sa = holding_info[d]
            port_map[d] = sa * (row["close"] / bp) if bp > 0 else sa
        else:
            port_map[d] = last_settled
        if d in sell_settled:
            last_settled = sell_settled[d]

    dates12  = pd.to_datetime(df12["date"])
    closes12 = df12["close"].values
    port12   = [port_map.get(d, last_settled) for d in df12["date"].values]

    cycles12 = [c for c in completed_cycles
                if c["sell_date"] >= cutoff or c["buy_date"] >= cutoff]

    # ── matplotlib 차트 ────────────────────────────────────────
    BG, FG, GRID = "#0d0d1a", "#c0c0d8", "#1e1e36"

    fig, ax1 = plt.subplots(figsize=(13, 5.4), facecolor=BG)
    ax1.set_facecolor(BG)
    for sp in ax1.spines.values():
        sp.set_color(GRID)
    ax1.grid(axis="y", color=GRID, linewidth=0.5, linestyle="--")
    ax1.tick_params(colors=FG, labelsize=9)

    # TQQQ 주가 (좌측 y)
    ax1.plot(dates12, closes12, color="#8899cc", linewidth=1.3, label=f"{ticker} 주가")
    ax1.set_ylabel(f"{ticker} ($)", color="#8899cc", fontsize=10)
    ax1.yaxis.label.set_color("#8899cc")
    ax1.tick_params(axis="y", colors="#8899cc")

    # 포트폴리오 자산 (우측 y)
    ax2 = ax1.twinx()
    ax2.set_facecolor(BG)
    for sp in ax2.spines.values():
        sp.set_color(GRID)
    ax2.plot(dates12, port12, color="#ffd700", linewidth=2.0, label="자산가치")
    ax2.set_ylabel("자산가치 ($)", color="#ffd700", fontsize=10)
    ax2.tick_params(colors="#ffd700", labelsize=9)
    ax2.yaxis.label.set_color("#ffd700")

    # 보유기간 음영
    for c in cycles12:
        bd  = max(c["buy_date"],  cutoff)
        sd  = c["sell_date"]
        clr = "#00bb44" if c["cycle_return_pct"] >= 0 else "#bb2200"
        try:
            ax1.axvspan(pd.Timestamp(bd), pd.Timestamp(sd),
                        alpha=0.22, color=clr, linewidth=0)
        except Exception:
            pass

    ax1.set_title(
        f"{ticker}  RSI(2)+F&G  최근 12개월  (보고서 기준 자산: 누적 복리)",
        color=FG, fontsize=11, pad=8,
    )
    ax1.xaxis.set_major_formatter(mdates.DateFormatter("%Y-%m"))
    ax1.xaxis.set_major_locator(mdates.MonthLocator(interval=1))
    plt.setp(ax1.xaxis.get_majorticklabels(), rotation=28, ha="right", fontsize=9)

    # 합체 범례
    h1, l1 = ax1.get_legend_handles_labels()
    h2, l2 = ax2.get_legend_handles_labels()
    extra = [
        mpatches.Patch(color="#00bb44", alpha=0.55, label="보유기간 (수익)"),
        mpatches.Patch(color="#bb2200", alpha=0.55, label="보유기간 (손실)"),
    ]
    ax1.legend(h1 + h2 + extra, l1 + l2 + [e.get_label() for e in extra],
               loc="upper left", fontsize=9, framealpha=0.3,
               facecolor=BG, edgecolor=GRID, labelcolor=FG)

    fig.tight_layout(pad=0.9)
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=150, facecolor=BG, bbox_inches="tight")
    plt.close(fig)
    buf.seek(0)
    return buf.read()


def _slide_backtest_chart(prs, ticker: str, df: pd.DataFrame,
                          completed_cycles: list):
    """슬라이드: 최근 12개월 백테스트 차트 이미지."""
    import io

    chart_png = _make_12m_chart(df, completed_cycles, ticker)

    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    _set_bg(s)
    _add_rect(s, 0, 0, SLIDE_W, Inches(0.08), ACCENT_CYAN)

    _add_text(s, f"{ticker}  최근 12개월  백테스트 차트",
              Inches(0.4), Inches(0.12), Inches(10), Inches(0.55),
              size=22, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    buf = io.BytesIO(chart_png)
    s.shapes.add_picture(buf, Inches(0.2), Inches(0.72),
                         width=Inches(12.93), height=Inches(6.38))

    _add_text(s,
              "※ 자산가치: 누적 복리 ($100K 시작, 연도 리셋 없음) · 수수료/세금 미포함",
              Inches(0.4), Inches(7.12), Inches(12.53), Inches(0.32),
              size=10, bold=False, color=TEXT_DIM, align=PP_ALIGN.RIGHT)


def _slide_disclaimer(prs):
    """슬라이드 6: 면책고지"""
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    _set_bg(s)
    _add_rect(s, 0, 0, SLIDE_W, Inches(0.08), ACCENT_GOLD)
    _add_rect(s, 0, SLIDE_H - Inches(0.08), SLIDE_W, Inches(0.08), ACCENT_GOLD)

    _add_text(s, "⚠  투자 면책 고지  /  Disclaimer",
              Inches(0.5), Inches(0.2), Inches(12), Inches(0.7),
              size=30, bold=True, color=ACCENT_GOLD, align=PP_ALIGN.CENTER)
    _add_divider(s, Inches(0.95), color=ACCENT_GOLD)

    disclaimers = [
        "본 유튜브 채널 및 관련 자료는 오직 교육 및 정보 제공 목적으로만 제작되었습니다.",
        "",
        "이 콘텐츠는 투자 권유, 투자 자문, 또는 특정 금융상품의 매수·매도를 권고하는 것이 아닙니다.",
        "모든 투자 결정은 시청자 본인의 판단과 책임 하에 이루어져야 하며,",
        "과거의 수익률이 미래의 수익률을 보장하지 않습니다.",
        "",
        "RSI(2) 전략은 레버리지 ETF(TQQQ, SOXL 등)를 포함하며,",
        "레버리지 ETF는 단기 변동성이 매우 크고 장기 보유 시 자산 감소 위험이 있습니다.",
        "투자 전 반드시 해당 상품의 투자설명서와 위험 고지를 숙지하시기 바랍니다.",
        "",
        "본 채널은 투자자문업 등록 없이 운영되며, 유료 투자 신호 서비스를 제공하지 않습니다.",
        "콘텐츠에 포함된 알고리즘 백테스트 결과는 수수료, 세금, 슬리피지를 일부 반영하나",
        "실제 투자 결과와 다를 수 있습니다.",
        "",
        "BATA Channel  ·  RSI(2) + Fear & Greed Algorithm  ·  For Educational Use Only",
    ]

    ty = Inches(1.1)
    for line in disclaimers:
        if not line:
            ty += Inches(0.18)
            continue
        _add_text(s, line, Inches(1.0), ty, Inches(11.3), Inches(0.42),
                  size=14, bold=False, color=GRAY_LIGHT)
        ty += Inches(0.38)


# ─────────────────────────────────────────────────────────────
# 메인 PPT 빌더
# ─────────────────────────────────────────────────────────────

def cfg_for(ticker: str) -> dict:
    return TICKER_CONFIG.get(ticker, TICKER_CONFIG["TQQQ"])


def build_cycle_ppt(
    ticker: str,
    cycle: dict,
    df: pd.DataFrame,
    ppt_type: str = "매도종료",
    completed_cycles: list | None = None,
) -> Path:
    """
    사이클 PPT 생성

    Parameters
    ----------
    ticker           : "TQQQ" or "SOXL"
    cycle            : extract_cycles() 반환값의 사이클 dict
    df               : simulate_fng() 반환 전체 DataFrame
    ppt_type         : "매수시작" | "매도종료" | "Breaking"
    completed_cycles : 완료된 전체 사이클 리스트 (연도별 복리 계산에 사용)
    """
    # 연도별 YY-N 번호 및 보고서 전용 투자금 계산
    year_info = _compute_year_info(cycle, completed_cycles or [cycle])
    # 사이클 기간 df 추출 (메모리 절약: 전체 df 는 이 함수 밖에서 del)
    buy_date  = cycle["buy_date"]
    sell_date = cycle.get("sell_date", "")

    cycle_mask = df["date"] >= buy_date
    if sell_date and sell_date != "진행중":
        cycle_mask &= df["date"] <= sell_date
    df_cycle = df[cycle_mask].copy()

    # 매수 전 7거래일 맥락 포함한 df (슬라이드 3, 5 용)
    buy_pos = df[df["date"] == buy_date].index
    if len(buy_pos):
        ctx_start = max(0, df.index.get_loc(buy_pos[0]) - 7)
    else:
        ctx_start = 0
    if sell_date and sell_date != "진행중":
        sell_pos = df[df["date"] == sell_date].index
        if len(sell_pos):
            ctx_end = min(len(df), df.index.get_loc(sell_pos[0]) + 2)
        else:
            ctx_end = len(df)
    else:
        ctx_end = len(df)
    df_context = df.iloc[ctx_start: ctx_end].copy()

    # PPT 생성
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    _slide_cover(prs, ticker, cycle, ppt_type, year_info)
    _slide_summary(prs, ticker, cycle, cfg_for(ticker), year_info, ppt_type)
    _slide_entry(prs, ticker, cycle, df_context)
    _slide_holding(prs, ticker, cycle, df_cycle)
    if sell_date and sell_date != "진행중":
        _slide_exit(prs, ticker, cycle, df_context)
    _slide_recent_cycles(prs, ticker, completed_cycles or [cycle])
    _slide_reference(prs, ticker, completed_cycles or [cycle])
    _slide_backtest_chart(prs, ticker, df, completed_cycles or [cycle])
    _slide_disclaimer(prs)

    # 파일명: TQQQ_Cycle05_매도종료_20260616.pptx
    today_str = datetime.now().strftime("%Y%m%d")
    fname     = f"{ticker}_Cycle{cycle['cycle_no']:02d}_{ppt_type}_{today_str}.pptx"
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    out_path  = OUT_DIR / fname
    prs.save(str(out_path))
    print(f"[PPT] 저장 완료: {out_path}")
    return out_path


# ─────────────────────────────────────────────────────────────
# 이메일 발송
# ─────────────────────────────────────────────────────────────

def send_email(ppt_path: Path, cycle: dict, ticker: str, ppt_type: str):
    app_password = os.environ.get("APP_PASS", "").replace(" ", "")
    if not app_password:
        print("⚠ APP_PASS 환경변수가 없습니다.")
        print("  APP_PASS=xxxx python make_cycle_ppt.py 형태로 실행하거나")
        print("  발급: https://myaccount.google.com/apppasswords")
        import getpass
        app_password = getpass.getpass("Gmail 앱 비밀번호 (16자리): ").replace(" ", "")

    ret_sign = "+" if cycle["cycle_return_pct"] >= 0 else ""
    subject  = (
        f"[BATA] {ticker} Cycle#{cycle['cycle_no']:02d} {ppt_type} "
        f"| 매수 {cycle['buy_date']} → 매도 {cycle.get('sell_date','-')} "
        f"| 수익률 {ret_sign}{cycle['cycle_return_pct']:.2f}%"
    )
    body = f"""안녕하세요,

{ticker} RSI(2)+FnG 전략 {ppt_type} 보고서를 첨부합니다.

■ 사이클 #{cycle['cycle_no']:02d} 요약
  매수일  : {cycle['buy_date']}  (종가 ${cycle['buy_price']:.2f}, RSI {cycle.get('buy_rsi') or '-'}, F&G {cycle.get('buy_fng') or '-'})
  매도일  : {cycle.get('sell_date', '-')}  (종가 ${cycle.get('sell_price', '-')}, RSI {cycle.get('sell_rsi') or '-'}, F&G {cycle.get('sell_fng') or '-'})
  보유기간: {cycle.get('days_held', '-')}일
  수익률  : {ret_sign}{cycle['cycle_return_pct']:.2f}%  (${abs(cycle.get('cycle_return_amount', 0)):,.0f})

⚠ 본 자료는 교육 목적이며 투자 권유가 아닙니다.

BATA · RSI(2) Algorithm Channel
"""

    msg = MIMEMultipart()
    msg["From"]    = SENDER
    msg["To"]      = RECEIVER
    msg["Subject"] = subject
    msg.attach(MIMEText(body, "plain", "utf-8"))

    # 첨부 파일명: 한글 → ASCII 변환 (MIME 헤더 인코딩 문제 방지)
    _type_ascii = {"매수시작": "BuyStart", "매도종료": "SellEnd", "Breaking": "Breaking"}
    email_filename = ppt_path.name
    for ko, en in _type_ascii.items():
        email_filename = email_filename.replace(ko, en)

    with open(ppt_path, "rb") as f:
        part = MIMEBase("application",
                        "vnd.openxmlformats-officedocument.presentationml.presentation")
        part.set_payload(f.read())
    encoders.encode_base64(part)
    part.add_header("Content-Disposition", f'attachment; filename="{email_filename}"')
    msg.attach(part)

    print(f"\n[메일] {RECEIVER} 으로 발송 중...")
    with smtplib.SMTP_SSL("smtp.gmail.com", 465) as server:
        server.login(SENDER, app_password)
        server.send_message(msg)
    print(f"[메일] 발송 완료! 제목: {subject}")


# ─────────────────────────────────────────────────────────────
# 메인
# ─────────────────────────────────────────────────────────────

def main():
    ticker = sys.argv[1].upper() if len(sys.argv) > 1 else "TQQQ"
    if ticker not in TICKER_CONFIG:
        print(f"[오류] 지원 종목: {list(TICKER_CONFIG)}")
        sys.exit(1)

    cfg = TICKER_CONFIG[ticker]

    # ── 1. 가격 데이터 ─────────────────────────────────────────
    print(f"[데이터] {ticker} 주가 다운로드 중... ({START_DATE} ~ today)")
    raw = yf.download(ticker, start=START_DATE, auto_adjust=True, progress=False)
    if raw.empty:
        print("[오류] yfinance 데이터 없음")
        sys.exit(1)
    close = raw["Close"].squeeze().dropna()
    close.index = pd.to_datetime(close.index).tz_localize(None)
    print(f"[데이터] {len(close)}거래일 로드 완료 ({close.index[0].date()} ~ {close.index[-1].date()})")

    # ── 2. F&G 데이터 ─────────────────────────────────────────
    print("[데이터] Fear & Greed 로드 중...")
    fng = load_fng_history()
    fng.index = pd.to_datetime(fng.index).tz_localize(None)

    # ── 3. 시뮬레이션 ─────────────────────────────────────────
    print("[시뮬] 백테스트 실행 중...")
    df = simulate_fng(
        close, fng,
        period=cfg["period"],
        buy_below=cfg["buy_below"],
        sell_above=cfg["sell_above"],
        fear_max=FEAR_MAX,
        greed_min=GREED_MIN,
    )
    print(f"[시뮬] 완료 — 총 {len(df)}행")

    # ── 4. 사이클 추출 ─────────────────────────────────────────
    completed, open_cycle = extract_cycles(df)
    print(f"[사이클] 완료 {len(completed)}개  |  진행중: {'있음' if open_cycle else '없음'}")

    if not completed and open_cycle is None:
        print("[오류] 사이클 없음")
        sys.exit(1)

    if completed:
        last_cycle = completed[-1]
        print(f"[사이클] 마지막 완료 사이클 #{last_cycle['cycle_no']:02d}: "
              f"{last_cycle['buy_date']} → {last_cycle['sell_date']}  "
              f"({'+' if last_cycle['cycle_return_pct'] >= 0 else ''}{last_cycle['cycle_return_pct']:.2f}%)")
    if open_cycle:
        print(f"[사이클] 진행중 사이클 #{open_cycle['cycle_no']:02d}: "
              f"{open_cycle['buy_date']} → 진행중  "
              f"(현재 {'+' if open_cycle['cycle_return_pct'] >= 0 else ''}{open_cycle['cycle_return_pct']:.2f}%)")

    # ── 5. PPT 생성 + 이메일 발송 + Google Drive 업로드 ────────
    # sys.argv[2] 로 타입 지정 가능: 매수시작 | 매도종료 | both (기본 both)
    requested = sys.argv[2] if len(sys.argv) > 2 else "both"
    types_to_send = (
        ["매수시작", "매도종료"] if requested == "both"
        else [requested]
    )

    for ppt_type in types_to_send:
        # 매수시작: 현재 진행중 사이클 우선 (없으면 마지막 완료 사이클)
        if ppt_type == "매수시작" and open_cycle is not None:
            target_cycle = open_cycle
            all_cycles   = completed + [open_cycle]
        elif completed:
            target_cycle = completed[-1]
            all_cycles   = completed
        else:
            print(f"[경고] {ppt_type} 용 완료 사이클 없음 — 스킵")
            continue

        print(f"[PPT] {ppt_type} 대상: Cycle#{target_cycle['cycle_no']:02d} "
              f"({target_cycle['buy_date']} → {target_cycle.get('sell_date', '-')})")

        ppt_path = build_cycle_ppt(
            ticker, target_cycle, df,
            ppt_type=ppt_type,
            completed_cycles=all_cycles,
        )
        send_email(ppt_path, target_cycle, ticker, ppt_type=ppt_type)

        # Google Drive 업로드
        try:
            from upload_to_gdrive import upload_ppt
            drive_result = upload_ppt(str(ppt_path))
            print(f"[Drive] 업로드 완료: {drive_result.get('webViewLink', '링크 없음')}")
        except Exception as e:
            print(f"[Drive] ⚠ 업로드 실패: {e}")

    # df 해제 (메모리)
    del df, close, fng, raw

    print("\n[완료] 모든 작업 종료")


if __name__ == "__main__":
    main()
